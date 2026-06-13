"""Tests for PPTX L2 Structure Detection and Content Fill.

Tests for:
- Structure parsing
- Placeholder detection
- Content mapping
- Style preservation
"""
import io
import pytest
from uuid import uuid4
from unittest.mock import AsyncMock, MagicMock, patch

from pptx import Presentation
from pptx.util import Inches, Pt

# Import the PPTXExporter and related classes
from app.domains.export.service import (
    PPTXExporter,
    PlaceholderInfo,
    SlideStructure,
    ContentMapping,
)


class TestPlaceholderInfo:
    """Tests for PlaceholderInfo dataclass."""

    def test_placeholder_info_creation(self):
        """Test creating a PlaceholderInfo object."""
        placeholder = PlaceholderInfo(
            variable_name="title",
            shape_id=1,
            shape_name="Title 1",
            placeholder_type="title",
            text="{{title}}",
            position=(0, 0),
        )

        assert placeholder.variable_name == "title"
        assert placeholder.shape_id == 1
        assert placeholder.shape_name == "Title 1"
        assert placeholder.placeholder_type == "title"
        assert placeholder.position == (0, 0)


class TestSlideStructure:
    """Tests for SlideStructure dataclass."""

    def test_slide_structure_creation(self):
        """Test creating a SlideStructure object."""
        structure = SlideStructure(
            slide_index=0,
            layout_name="Title Slide",
            layout_type="title",
        )

        assert structure.slide_index == 0
        assert structure.layout_name == "Title Slide"
        assert structure.layout_type == "title"
        assert structure.placeholders == []
        assert structure.shapes_with_placeholders == []


class TestPPTXExporterStructureParsing:
    """Tests for PPTX structure parsing methods."""

    @pytest.fixture
    def exporter(self):
        """Create a PPTXExporter instance."""
        db = AsyncMock()
        tenant_id = uuid4()
        return PPTXExporter(db, tenant_id)

    @pytest.fixture
    def sample_presentation(self):
        """Create a sample presentation with placeholders."""
        prs = Presentation()
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "{{presentation_title}}"

        # Add content slide
        content_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(content_layout)

        # Find body placeholder
        for shape in slide2.shapes:
            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                if "BODY" in str(shape.placeholder_format.type):
                    tf = shape.text_frame
                    tf.clear()
                    tf.paragraphs[0].text = "{{executive_summary}}"
                    break

        return prs

    def test_contains_placeholder(self, exporter):
        """Test placeholder detection in text."""
        assert exporter._contains_placeholder("{{title}}") is True
        assert exporter._contains_placeholder("Hello {{name}}") is True
        assert exporter._contains_placeholder("客户：{{客户名称}}") is True
        assert exporter._contains_placeholder("编号：{{project.code-1}}") is True
        assert exporter._contains_placeholder("No placeholder") is False
        assert exporter._contains_placeholder("{{123invalid}}") is False  # Can't start with number
        assert exporter._contains_placeholder("{{invalid name}}") is False

    def test_extract_placeholders_from_text(self, exporter):
        """Test extracting variable names from placeholder text."""
        result = exporter._extract_placeholders_from_text("{{title}} and {{subtitle}}")
        assert result == ["title", "subtitle"]

        result = exporter._extract_placeholders_from_text("Single {{var}}")
        assert result == ["var"]

        result = exporter._extract_placeholders_from_text("{{客户名称}} / {{业务范围}} / {{project.code-1}}")
        assert result == ["客户名称", "业务范围", "project.code-1"]

        result = exporter._extract_placeholders_from_text("No placeholders")
        assert result == []

    def test_detect_slide_layouts(self, exporter, sample_presentation):
        """Test detecting slide layout types."""
        layouts = exporter._detect_slide_layouts(sample_presentation)

        assert len(layouts) == 2
        # First slide should be title layout
        assert layouts[0]["layout_type"] == "title"
        assert layouts[0]["has_title"] is True
        # Second slide should have some structure (title layout if only title placeholder detected)
        assert layouts[1]["layout_type"] in ("title", "content", "two_column")

    def test_extract_placeholders(self, exporter, sample_presentation):
        """Test extracting placeholders from slides."""
        slide1 = sample_presentation.slides[0]
        placeholders = exporter._extract_placeholders(slide1)

        assert len(placeholders) >= 1
        var_names = [p.variable_name for p in placeholders]
        assert "presentation_title" in var_names

    def test_parse_content_sections(self, exporter):
        """Test parsing markdown content into sections."""
        content = """# Introduction
This is the introduction text.

# Executive Summary
This is the executive summary.

## Details
These are the details.
"""
        sections = exporter._parse_content_sections(content)

        assert "Introduction" in sections
        assert "Executive Summary" in sections
        assert "Details" in sections
        assert "This is the introduction text." in sections["Introduction"]

    def test_find_content_for_placeholder_title(self, exporter):
        """Test finding content for title placeholders."""
        sections = {
            "Introduction": "Intro content",
            "Executive Summary": "Summary content",
        }

        result = exporter._find_content_for_placeholder("title", "title", sections)
        # For title type, it returns the first section key (section header name)
        assert result == "Introduction"

    def test_find_content_for_placeholder_body(self, exporter):
        """Test finding content for body placeholders."""
        sections = {
            "default": "Default section content",
            "Section 1": "Section 1 content",
        }

        result = exporter._find_content_for_placeholder("body", "body", sections)
        assert result == "Default section content"

    def test_find_content_for_placeholder_by_var_name(self, exporter):
        """Test finding content by variable name matching."""
        sections = {
            "Executive Summary": "Summary content",
            "Introduction": "Intro content",
        }

        result = exporter._find_content_for_placeholder("executive_summary", "body", sections)
        assert result == "Summary content"


class TestPPTXExporterContentFill:
    """Tests for PPTX content fill methods."""

    @pytest.fixture
    def exporter(self):
        """Create a PPTXExporter instance."""
        db = AsyncMock()
        tenant_id = uuid4()
        return PPTXExporter(db, tenant_id)

    @pytest.fixture
    def presentation_with_placeholder(self):
        """Create a presentation with a placeholder."""
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "{{document_title}}"
        return prs

    def test_substitute_variables(self, exporter):
        """Test variable substitution."""
        content = "Hello {{name}}, your email is {{email}}"
        variables = {"name": "John", "email": "john@example.com"}

        result = exporter._substitute_variables(content, variables)
        assert result == "Hello John, your email is john@example.com"

    def test_substitute_chinese_variables_and_preserve_missing_values(self, exporter):
        """Test Chinese variable substitution without hiding unresolved placeholders."""
        content = "客户：{{客户名称}}，范围：{{业务范围}}，缺失：{{未配置变量}}"
        variables = {"客户名称": "远大客户", "业务范围": "仓储数字化"}

        result = exporter._substitute_variables(content, variables)
        assert result == "客户：远大客户，范围：仓储数字化，缺失：{{未配置变量}}"

    def test_fill_slide_placeholders(self, exporter, presentation_with_placeholder):
        """Test filling placeholders in a slide."""
        slide = presentation_with_placeholder.slides[0]
        # content_map keys should be variable names without {{}} - _substitute_variables adds them
        content_map = {"document_title": "My Document Title"}

        exporter._fill_slide_placeholders(slide, content_map)

        # Check that the placeholder was replaced
        title_shape = slide.shapes.title
        assert title_shape.text == "My Document Title"

    def test_set_text_frame_preserving_style(self, exporter):
        """Test setting text frame content while preserving styles."""
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title

        # Set initial text
        title_shape.text = "{{title}}"

        # Fill with content
        new_text = "New Title Text"
        exporter._set_text_frame_preserving_style(title_shape.text_frame, new_text)

        assert title_shape.text == new_text

    def test_map_content_to_slides(self, exporter):
        """Test mapping document content to slide structures."""
        content = """# Title Section
Title content here.

# Summary Section
Summary content here.
"""
        slide_map = [
            SlideStructure(
                slide_index=0,
                layout_name="Title Slide",
                layout_type="title",
                placeholders=[
                    PlaceholderInfo(
                        variable_name="title",
                        shape_id=1,
                        shape_name="Title 1",
                        placeholder_type="title",
                        text="{{title}}",
                        position=(0, 0),
                    )
                ],
            ),
            SlideStructure(
                slide_index=1,
                layout_name="Content Slide",
                layout_type="content",
                placeholders=[
                    PlaceholderInfo(
                        variable_name="summary",
                        shape_id=2,
                        shape_name="Body 1",
                        placeholder_type="body",
                        text="{{summary}}",
                        position=(0, 0),
                    )
                ],
            ),
        ]

        content_map = exporter._map_content_to_slides(content, slide_map)

        assert 0 in content_map
        assert "title" in content_map[0]
        assert 1 in content_map
        assert "summary" in content_map[1]


class TestPPTXExporterIntegration:
    """Integration tests for PPTXExporter."""

    @pytest.fixture
    def exporter(self):
        """Create a PPTXExporter instance."""
        db = AsyncMock()
        tenant_id = uuid4()
        return PPTXExporter(db, tenant_id)

    def test_get_presentation_structure(self, exporter):
        """Test getting complete presentation structure."""
        prs = Presentation()
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "{{title}}"

        # Add content slide
        content_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(content_layout)

        structure = exporter.get_presentation_structure(prs)

        assert structure["slide_count"] == 2
        assert len(structure["slides"]) == 2
        assert len(structure["layouts"]) == 2
        # First slide should have title layout
        assert structure["slides"][0]["layout_type"] == "title"
        # Should detect the placeholder
        assert structure["slides"][0]["placeholder_count"] >= 1

    def test_export_with_template_variables(self, exporter):
        """Test export with template variable substitution."""
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "{{project_name}} - {{version}}"

        # Apply variable substitution
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                tf = shape.text_frame
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        if exporter._contains_placeholder(run.text):
                            run.text = exporter._substitute_variables(
                                run.text, {"project_name": "MyProject", "version": "1.0"}
                            )

        title_shape = slide.shapes.title
        assert "MyProject" in title_shape.text
        assert "1.0" in title_shape.text

    def test_strict_placeholder_pattern(self, exporter):
        """Test that only valid placeholders are detected."""
        # Valid placeholders
        assert exporter._contains_placeholder("{{valid_var}}") is True
        assert exporter._contains_placeholder("{{var123}}") is True
        assert exporter._contains_placeholder("{{_under}}") is True
        assert exporter._contains_placeholder("{{camelCase}}") is True
        assert exporter._contains_placeholder("{{project.code-1}}") is True
        assert exporter._contains_placeholder("{{客户名称}}") is True

        # Invalid placeholders (can't start with number)
        assert exporter._contains_placeholder("{{123invalid}}") is False
        assert exporter._contains_placeholder("{{invalid name}}") is False
        assert exporter._contains_placeholder("{{invalid/name}}") is False


class TestContentMapping:
    """Tests for ContentMapping dataclass."""

    def test_content_mapping_creation(self):
        """Test creating a ContentMapping object."""
        mapping = ContentMapping(
            slide_index=0,
            placeholder_type="title",
            variable_name="title",
            content="My Title",
            style_info={"font_size": 24},
        )

        assert mapping.slide_index == 0
        assert mapping.placeholder_type == "title"
        assert mapping.variable_name == "title"
        assert mapping.content == "My Title"
        assert mapping.style_info["font_size"] == 24
