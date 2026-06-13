"""Tests for project-level delivery package exports."""

import io
import os
from datetime import datetime, timezone
from uuid import uuid4

os.environ["DATABASE_URL"] = "sqlite+aiosqlite:///./test-project-export-package.db"
os.environ["REDIS_URL"] = "redis://localhost:6379/0"
os.environ["ARQ_REDIS_URL"] = "redis://localhost:6379/1"
os.environ["JWT_SECRET_KEY"] = "test-project-export-package-secret"

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from sqlalchemy.ext.asyncio import async_sessionmaker, create_async_engine

from app.core.settings import settings

settings.DATABASE_URL = os.environ["DATABASE_URL"]
settings.REDIS_URL = os.environ["REDIS_URL"]
settings.ARQ_REDIS_URL = os.environ["ARQ_REDIS_URL"]

import app.db.init_schema  # noqa: F401 - registers sqlite compilers for UUID/JSONB
from app.db.base import Base
from app.db.init_schema import deduplicate_indexes
from app.domains.documents.models import Document, DocumentStatus, DocumentType
from app.domains.export.models import ExportArtifact, ExportJob, ExportStatus, ExportType
from app.domains.export.schemas import ProjectPackageExportRequest
from app.domains.export.service import ExportService
from app.models.identity import Tenant, User
from app.models.projects import Project, ProjectMember
from app.services.storage import StorageHandle


@pytest.fixture
async def db_session():
    """Create an isolated async SQLite database for export package tests."""
    deduplicate_indexes()
    engine = create_async_engine("sqlite+aiosqlite:///:memory:")
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)

    Session = async_sessionmaker(engine, expire_on_commit=False)
    async with Session() as session:
        yield session

    await engine.dispose()


class CapturingStorage:
    """Storage fake that keeps the last uploaded artifact in memory."""

    def __init__(self):
        self.uploads = []

    async def upload(self, tenant_id, project_id, filename, content, content_type):
        self.uploads.append(
            {
                "tenant_id": tenant_id,
                "project_id": project_id,
                "filename": filename,
                "content": content,
                "content_type": content_type,
            }
        )
        return StorageHandle(
            path=f"{tenant_id}/{project_id}/{filename}",
            filename=filename,
            content_type=content_type,
            size=len(content),
            hash="b" * 64,
            storage_backend="test",
        )


def _document(*, tenant_id, project_id, user_id, doc_type, title, status, content=None, metadata=None):
    return Document(
        id=uuid4(),
        tenant_id=tenant_id,
        project_id=project_id,
        doc_type=doc_type,
        title=title,
        content=content or f"# {title}\n\n{{{{client_name}}}} delivery content.",
        status=status,
        version=1,
        created_by=user_id,
        metadata_json=metadata or {},
    )


async def _seed_project(db_session):
    tenant_id = uuid4()
    user_id = uuid4()
    project_id = uuid4()

    tenant = Tenant(id=tenant_id, name="Export Tenant", slug="export-tenant")
    user = User(
        id=user_id,
        tenant_id=tenant_id,
        email="owner@example.com",
        hashed_password="test",
        full_name="Project Owner",
    )
    project = Project(
        id=project_id,
        tenant_id=tenant_id,
        owner_id=user_id,
        name="智能制造交付项目",
        slug="smart-manufacturing-delivery",
    )
    member = ProjectMember(project_id=project_id, user_id=user_id)

    urs = _document(
        tenant_id=tenant_id,
        project_id=project_id,
        user_id=user_id,
        doc_type=DocumentType.URS.value,
        title="URS 业务目标",
        status=DocumentStatus.PUBLISHED.value,
    )
    prd = _document(
        tenant_id=tenant_id,
        project_id=project_id,
        user_id=user_id,
        doc_type=DocumentType.PRD.value,
        title="PRD 产品能力",
        status=DocumentStatus.PUBLISHED.value,
    )
    draft = _document(
        tenant_id=tenant_id,
        project_id=project_id,
        user_id=user_id,
        doc_type=DocumentType.TEST_CASE.value,
        title="测试用例草稿",
        status=DocumentStatus.DRAFT.value,
    )

    db_session.add_all([tenant, user, project, member, urs, prd, draft])
    await db_session.flush()
    return tenant_id, user_id, project_id, urs, prd, draft


@pytest.mark.asyncio
async def test_project_package_exports_selected_documents_in_delivery_order(db_session, monkeypatch):
    tenant_id, user_id, project_id, urs, prd, draft = await _seed_project(db_session)
    storage = CapturingStorage()
    monkeypatch.setattr("app.domains.export.service.get_storage_provider", lambda: storage)

    service = ExportService(db_session)
    job = await service.export_project_package(
        project_id=project_id,
        tenant_id=tenant_id,
        document_ids=[prd.id, urs.id],
        title="客户交付包",
        variables={"client_name": "远大客户"},
        created_by=user_id,
    )

    assert job.export_type == ExportType.PROJECT_PACKAGE.value
    assert job.status == ExportStatus.COMPLETED.value
    assert job.document_id is None
    assert len(job.artifacts) == 1
    assert job.artifacts[0].filename.endswith(".project-package.md")

    uploaded = storage.uploads[0]
    content = uploaded["content"].decode("utf-8")
    assert uploaded["content_type"] == "text/markdown"
    assert "# 客户交付包" in content
    assert "## 交付清单" in content
    assert content.index("URS 业务目标") < content.index("PRD 产品能力")
    assert "远大客户 delivery content" in content
    assert "测试用例草稿" not in content


def test_project_package_request_preserves_production_options():
    project_id = uuid4()
    document_id = uuid4()

    request = ProjectPackageExportRequest(
        project_id=project_id,
        document_ids=[document_id],
        formats=["markdown", "word", "pptx"],
        include_audit=True,
        watermark="Client review copy",
    )

    assert request.formats == ["markdown", "word", "pptx"]
    assert request.include_audit is True
    assert request.watermark == "Client review copy"


@pytest.mark.asyncio
async def test_project_package_generates_selected_format_artifacts_with_audit_and_watermark(db_session, monkeypatch):
    tenant_id, user_id, project_id, urs, prd, _ = await _seed_project(db_session)
    storage = CapturingStorage()
    monkeypatch.setattr("app.domains.export.service.get_storage_provider", lambda: storage)

    service = ExportService(db_session)
    job = await service.export_project_package(
        project_id=project_id,
        tenant_id=tenant_id,
        document_ids=[prd.id, urs.id],
        title="Production delivery pack",
        formats=["markdown", "word", "pptx"],
        include_audit=True,
        watermark="Client review copy",
        variables={"client_name": "ACME"},
        created_by=user_id,
    )

    assert job.status == ExportStatus.COMPLETED.value
    assert len(job.artifacts) == 3
    filenames = {artifact.filename for artifact in job.artifacts}
    assert "Production-delivery-pack.project-package.md" in filenames
    assert "Production-delivery-pack.project-package.docx" in filenames
    assert "Production-delivery-pack.project-package.pptx" in filenames

    uploads_by_filename = {upload["filename"]: upload for upload in storage.uploads}
    markdown = uploads_by_filename["Production-delivery-pack.project-package.md"]["content"].decode("utf-8")
    assert "Client review copy" in markdown
    assert "Audit checklist" in markdown
    assert uploads_by_filename["Production-delivery-pack.project-package.docx"]["content_type"] == (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
    assert uploads_by_filename["Production-delivery-pack.project-package.pptx"]["content_type"] == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


@pytest.mark.asyncio
async def test_project_package_substitutes_chinese_placeholders_across_formats(db_session, monkeypatch):
    tenant_id, user_id, project_id, urs, _, _ = await _seed_project(db_session)
    urs.content = "# 项目摘要\n客户：{{客户名称}}\n范围：{{业务范围}}\n缺失：{{未配置变量}}"
    storage = CapturingStorage()
    monkeypatch.setattr("app.domains.export.service.get_storage_provider", lambda: storage)

    service = ExportService(db_session)
    job = await service.export_project_package(
        project_id=project_id,
        tenant_id=tenant_id,
        document_ids=[urs.id],
        title="Chinese placeholder package",
        formats=["markdown", "word", "pptx"],
        variables={"客户名称": "远大客户", "业务范围": "仓储数字化"},
        created_by=user_id,
    )

    assert job.status == ExportStatus.COMPLETED.value
    uploads_by_filename = {upload["filename"]: upload for upload in storage.uploads}

    markdown = uploads_by_filename["Chinese-placeholder-package.project-package.md"]["content"].decode("utf-8")
    assert "客户：远大客户" in markdown
    assert "范围：仓储数字化" in markdown
    assert "{{未配置变量}}" in markdown

    docx = DocxDocument(io.BytesIO(uploads_by_filename["Chinese-placeholder-package.project-package.docx"]["content"]))
    docx_text = "\n".join(paragraph.text for paragraph in docx.paragraphs)
    assert "客户：远大客户" in docx_text
    assert "范围：仓储数字化" in docx_text
    assert "{{未配置变量}}" in docx_text

    pptx = Presentation(io.BytesIO(uploads_by_filename["Chinese-placeholder-package.project-package.pptx"]["content"]))
    pptx_text = "\n".join(
        shape.text
        for slide in pptx.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "客户：远大客户" in pptx_text
    assert "范围：仓储数字化" in pptx_text
    assert "{{未配置变量}}" in pptx_text


@pytest.mark.asyncio
async def test_project_package_blocks_placeholder_documents(db_session, monkeypatch):
    tenant_id, user_id, project_id, _, _, draft = await _seed_project(db_session)
    draft.metadata_json = {"generation_status": "placeholder"}
    draft.status = DocumentStatus.PUBLISHED.value
    storage = CapturingStorage()
    monkeypatch.setattr("app.domains.export.service.get_storage_provider", lambda: storage)

    service = ExportService(db_session)
    with pytest.raises(ValueError, match="placeholder documents"):
        await service.export_project_package(
            project_id=project_id,
            tenant_id=tenant_id,
            document_ids=[draft.id],
            created_by=user_id,
        )

    assert storage.uploads == []


@pytest.mark.asyncio
async def test_project_export_readiness_reports_missing_types_and_blockers(db_session):
    tenant_id, _, project_id, urs, prd, draft = await _seed_project(db_session)
    draft.metadata_json = {"generation_status": "placeholder"}
    draft.status = DocumentStatus.PUBLISHED.value

    service = ExportService(db_session)
    readiness = await service.get_project_export_readiness(
        project_id=project_id,
        tenant_id=tenant_id,
    )

    assert readiness.total_documents == 3
    assert readiness.exportable_documents == 2
    assert readiness.blocked_documents == 1
    assert readiness.can_export_production is False
    assert "brd" in readiness.missing_required_types
    assert "detailed_design" in readiness.missing_required_types
    assert any(item.doc_type == "test_case" and item.status == "blocked" for item in readiness.required_types)
    assert readiness.blockers[0].document_id == draft.id
    assert readiness.blockers[0].reason == "文档仍包含占位内容"


@pytest.mark.asyncio
async def test_project_release_evidence_summarizes_gate_artifacts_and_actions(db_session):
    tenant_id, user_id, project_id, _, _, draft = await _seed_project(db_session)
    draft.metadata_json = {"generation_status": "placeholder"}
    draft.status = DocumentStatus.PUBLISHED.value

    completed_job = ExportJob(
        tenant_id=tenant_id,
        project_id=project_id,
        document_id=None,
        template_id=None,
        export_type=ExportType.PROJECT_PACKAGE.value,
        status=ExportStatus.COMPLETED.value,
        output_path="exports/package.md",
        file_hash="a" * 64,
        created_by=user_id,
        completed_at=datetime.now(timezone.utc),
    )
    failed_job = ExportJob(
        tenant_id=tenant_id,
        project_id=project_id,
        document_id=None,
        template_id=None,
        export_type=ExportType.PROJECT_PACKAGE.value,
        status=ExportStatus.FAILED.value,
        error_message="PPTX template variable missing",
        created_by=user_id,
        completed_at=datetime.now(timezone.utc),
    )
    db_session.add_all([completed_job, failed_job])
    await db_session.flush()
    db_session.add_all([
        ExportArtifact(
            tenant_id=tenant_id,
            job_id=completed_job.id,
            filename="Delivery.project-package.md",
            content_type="text/markdown",
            file_size=2048,
            storage_path="exports/Delivery.project-package.md",
            file_hash="b" * 64,
        ),
        ExportArtifact(
            tenant_id=tenant_id,
            job_id=completed_job.id,
            filename="Delivery.project-package.docx",
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            file_size=4096,
            storage_path="exports/Delivery.project-package.docx",
            file_hash="c" * 64,
        ),
    ])
    await db_session.flush()

    evidence = await ExportService(db_session).get_project_release_evidence(
        project_id=project_id,
        tenant_id=tenant_id,
    )

    assert evidence.release_gate.status == "blocked"
    assert evidence.summary.completed_jobs == 1
    assert evidence.summary.failed_jobs == 1
    assert evidence.summary.artifact_count == 2
    assert evidence.summary.covered_formats == ["Markdown", "Word"]
    assert evidence.summary.missing_formats == ["PPTX"]
    assert evidence.latest_job.id == completed_job.id
    assert evidence.recent_artifacts[0].filename == "Delivery.project-package.docx"
    assert {item.code for item in evidence.risk_items} >= {
        "export_readiness_blocked",
        "failed_exports",
        "missing_format_coverage",
    }
    assert evidence.priority_actions[0].href == "/documents"
