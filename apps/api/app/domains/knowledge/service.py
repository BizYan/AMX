"""Knowledge Domain Service

Business logic for knowledge management, RAG, and lineage tracking.
"""

import asyncio
import hashlib
import re
from datetime import datetime, timezone
from typing import Any
from uuid import UUID

from sqlalchemy import select, func, or_, and_
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import selectinload

from app.core.settings import settings
from app.integrations.llm.gateway import LLMGateway
from app.domains.knowledge.models import (
    KnowledgeEntry,
    KnowledgeLink,
    ProvenanceRecord,
    KnowledgeVector,
    LineageRecord,
)
from app.domains.knowledge.schemas import (
    KnowledgeEntryCreate,
    KnowledgeEntryUpdate,
    KnowledgeEntryResponse,
    KnowledgeLinkCreate,
    KnowledgeLinkResponse,
    KnowledgeGraphEdge,
    KnowledgeGraphNode,
    KnowledgeGraphSummary,
    KnowledgeGraphWorkbenchResponse,
    ProvenanceRecordCreate,
    ProvenanceRecordResponse,
    LineageRecordCreate,
    LineageRecordResponse,
    PaginatedResponse,
    GraphNeighborResponse,
    GraphPathResponse,
    BulkEntryCreate,
    BulkEntryResponse,
)
from app.services.storage import get_storage_provider, StorageHandle
from app.services.vector_store_provider import PGVectorProvider
from app.services.search_provider import PostgresFTSProvider
from app.services.graph_store_provider import PostgresGraphProvider


class SourceTextExtractionError(ValueError):
    """Raised when source file text cannot be extracted for knowledge ingestion."""


class KnowledgeService:
    """Service for knowledge management operations.

    Handles CRUD operations for knowledge entries, linking, lineage,
    and integrates with vector store and search providers.
    """

    def __init__(self, db: AsyncSession, llm_gateway: LLMGateway | None = None):
        """Initialize knowledge service.

        Args:
            db: Async database session
            llm_gateway: Optional LLM gateway for embedding generation
        """
        self.db = db
        self.llm_gateway = llm_gateway
        self.vector_provider = PGVectorProvider(db)
        self.search_provider = PostgresFTSProvider(db)
        self.graph_provider = PostgresGraphProvider(db)

    def _entry_to_response(self, entry: KnowledgeEntry) -> KnowledgeEntryResponse:
        return KnowledgeEntryResponse(
            id=entry.id,
            tenant_id=entry.tenant_id,
            project_id=entry.project_id,
            source_file_id=entry.source_file_id,
            entry_type=entry.entry_type,
            content=entry.content,
            content_hash=entry.content_hash,
            vector_embedding=entry.vector_embedding,
            metadata=entry.metadata_json,
            sharing_scope=entry.sharing_scope,
            created_at=entry.created_at,
            updated_at=entry.updated_at,
            deleted_at=entry.deleted_at,
            created_by_id=entry.created_by_id,
            reviewed_by_id=entry.reviewed_by_id,
            reviewed_at=entry.reviewed_at,
        )

    def _entry_to_graph_node(self, entry: KnowledgeEntry) -> KnowledgeGraphNode:
        metadata = entry.metadata_json or {}
        title = str(
            metadata.get("title")
            or metadata.get("name")
            or metadata.get("heading")
            or entry.content[:42]
            or entry.id
        )
        summary = str(metadata.get("summary") or metadata.get("description") or entry.content[:240])
        source_label = str(
            metadata.get("sourceFileName")
            or metadata.get("source_file_name")
            or metadata.get("filename")
            or metadata.get("source")
            or "项目知识库"
        )
        raw_status = str(metadata.get("state") or metadata.get("status") or "").lower()
        if raw_status in {"gap", "missing", "todo"} or metadata.get("gap"):
            status_value = "gap"
        elif raw_status in {"conflict", "contradiction"} or metadata.get("conflict"):
            status_value = "conflict"
        else:
            status_value = "ready"

        confidence = metadata.get("confidence")
        try:
            confidence_value = float(confidence) if confidence is not None else None
        except (TypeError, ValueError):
            confidence_value = None

        return KnowledgeGraphNode(
            id=entry.id,
            type=entry.entry_type,
            title=title,
            summary=summary,
            project_id=entry.project_id,
            source_file_id=entry.source_file_id,
            source_label=source_label,
            status=status_value,
            confidence=confidence_value,
            metadata=metadata,
            created_at=entry.created_at,
            updated_at=entry.updated_at,
        )

    def _link_to_graph_edge(self, link: KnowledgeLink) -> KnowledgeGraphEdge:
        return KnowledgeGraphEdge(
            id=link.id,
            source=link.source_entry_id,
            target=link.target_entry_id,
            type=link.link_type,
            confidence=link.confidence,
            metadata=link.metadata_json or {},
            created_at=link.created_at,
        )

    async def create_entry(
        self,
        tenant_id: UUID,
        project_id: UUID,
        entry_type: str,
        content: str,
        source_file_id: UUID | None = None,
        metadata: dict[str, Any] | None = None,
        generate_embedding: bool = True,
        sharing_scope: str = "project",
        created_by_id: UUID | None = None,
    ) -> KnowledgeEntryResponse:
        """Create a new knowledge entry.

        Args:
            tenant_id: Tenant UUID
            project_id: Project UUID
            entry_type: Entry type (text, table, image, code)
            content: Content text
            source_file_id: Optional source file UUID
            metadata: Optional metadata dict
            generate_embedding: Whether to generate vector embedding

        Returns:
            Created KnowledgeEntryResponse
        """
        # Calculate content hash
        content_hash = hashlib.sha256(content.encode()).hexdigest()

        # Create entry
        entry = KnowledgeEntry(
            tenant_id=tenant_id,
            project_id=project_id,
            source_file_id=source_file_id,
            entry_type=entry_type,
            content=content,
            content_hash=content_hash,
            metadata_json=metadata,
            sharing_scope=sharing_scope,
            created_by_id=created_by_id,
        )
        self.db.add(entry)
        await self.db.flush()
        await self.db.refresh(entry)

        # Generate embedding if requested
        vector_embedding = None
        if generate_embedding and self.llm_gateway:
            try:
                embed_response = await self.llm_gateway.embed([content])
                if embed_response.embeddings:
                    vector_embedding = embed_response.embeddings[0]
                    entry.vector_embedding = vector_embedding

                    # Store in vector table
                    await self.vector_provider.upsert_vector(
                        entry.id,
                        vector_embedding,
                        {"vector_index": "default"},
                    )
            except Exception:
                # Log error but don't fail entry creation
                pass

        # Index for full-text search
        try:
            await self.search_provider.index_document(
                entry.id,
                content,
                metadata,
            )
        except Exception:
            # Log error but don't fail entry creation
            pass

        await self.db.flush()

        response = self._entry_to_response(entry)
        response.vector_embedding = vector_embedding
        return response

    async def get_entry(self, entry_id: UUID, tenant_id: UUID, user_id: UUID | None = None) -> KnowledgeEntryResponse | None:
        """Get a knowledge entry by ID.

        Args:
            entry_id: Entry UUID
            tenant_id: Tenant UUID for verification
            user_id: User UUID for sharing scope verification

        Returns:
            KnowledgeEntryResponse if found, None otherwise
        """
        result = await self.db.execute(
            select(KnowledgeEntry).where(
                KnowledgeEntry.id == entry_id,
                KnowledgeEntry.tenant_id == tenant_id,
                KnowledgeEntry.deleted_at.is_(None),
            )
        )
        entry = result.scalar_one_or_none()

        if not entry:
            return None

        # Check sharing scope - PRIVATE entries only visible to creator
        if entry.sharing_scope == "private" and entry.created_by_id != user_id:
            return None

        return self._entry_to_response(entry)

    async def list_entries(
        self,
        tenant_id: UUID,
        project_id: UUID | None = None,
        entry_type: str | None = None,
        source_file_id: UUID | None = None,
        page: int = 1,
        page_size: int = 20,
        user_id: UUID | None = None,
    ) -> PaginatedResponse[KnowledgeEntryResponse]:
        """List knowledge entries with pagination and filters.

        Args:
            tenant_id: Tenant UUID
            project_id: Optional project UUID filter
            entry_type: Optional entry type filter
            source_file_id: Optional source file UUID filter
            page: Page number (1-indexed)
            page_size: Items per page
            user_id: User UUID for sharing scope filter

        Returns:
            PaginatedResponse of KnowledgeEntryResponse
        """
        # Build query with sharing scope filter
        query = select(KnowledgeEntry).where(
            KnowledgeEntry.tenant_id == tenant_id,
            KnowledgeEntry.deleted_at.is_(None),
        )

        # Apply sharing scope filter based on user access
        # - PRIVATE: only creator
        # - PROJECT: project members
        # - TENANT: all tenant users
        # - GLOBAL: all users
        if user_id:
            query = query.where(
                or_(
                    KnowledgeEntry.sharing_scope.in_(["project", "tenant", "global"]),
                    and_(
                        KnowledgeEntry.sharing_scope == "private",
                        KnowledgeEntry.created_by_id == user_id,
                    ),
                )
            )

        if project_id:
            query = query.where(KnowledgeEntry.project_id == project_id)
        if entry_type:
            query = query.where(KnowledgeEntry.entry_type == entry_type)
        if source_file_id:
            query = query.where(KnowledgeEntry.source_file_id == source_file_id)

        # Count total
        count_query = select(func.count(KnowledgeEntry.id)).where(
            KnowledgeEntry.tenant_id == tenant_id,
            KnowledgeEntry.deleted_at.is_(None),
        )
        if user_id:
            count_query = count_query.where(
                or_(
                    KnowledgeEntry.sharing_scope.in_(["project", "tenant", "global"]),
                    and_(
                        KnowledgeEntry.sharing_scope == "private",
                        KnowledgeEntry.created_by_id == user_id,
                    ),
                )
            )
        if project_id:
            count_query = count_query.where(KnowledgeEntry.project_id == project_id)
        if entry_type:
            count_query = count_query.where(KnowledgeEntry.entry_type == entry_type)
        if source_file_id:
            count_query = count_query.where(KnowledgeEntry.source_file_id == source_file_id)

        count_result = await self.db.execute(count_query)
        total = count_result.scalar_one()

        # Get paginated results
        offset = (page - 1) * page_size
        query = query.offset(offset).limit(page_size).order_by(KnowledgeEntry.created_at.desc())

        result = await self.db.execute(query)
        entries = list(result.scalars().all())

        items = [self._entry_to_response(e) for e in entries]

        return PaginatedResponse(
            items=items,
            total=total,
            page=page,
            page_size=page_size,
            has_more=(offset + len(items)) < total,
        )

    async def update_entry(
        self,
        entry_id: UUID,
        tenant_id: UUID,
        updates: KnowledgeEntryUpdate,
    ) -> KnowledgeEntryResponse | None:
        """Update a knowledge entry.

        Args:
            entry_id: Entry UUID
            tenant_id: Tenant UUID for verification
            updates: Update data

        Returns:
            Updated KnowledgeEntryResponse if found, None otherwise
        """
        result = await self.db.execute(
            select(KnowledgeEntry).where(
                KnowledgeEntry.id == entry_id,
                KnowledgeEntry.tenant_id == tenant_id,
                KnowledgeEntry.deleted_at.is_(None),
            )
        )
        entry = result.scalar_one_or_none()

        if not entry:
            return None

        # Update fields
        if updates.content is not None:
            entry.content = updates.content
            entry.content_hash = hashlib.sha256(updates.content.encode()).hexdigest()

            # Re-generate embedding if content changed
            if self.llm_gateway:
                try:
                    embed_response = await self.llm_gateway.embed([updates.content])
                    if embed_response.embeddings:
                        entry.vector_embedding = embed_response.embeddings[0]
                        await self.vector_provider.upsert_vector(
                            entry_id,
                            embed_response.embeddings[0],
                            {"vector_index": "default"},
                        )
                except Exception:
                    pass

            # Re-index for FTS
            try:
                await self.search_provider.reindex_entry(entry_id)
            except Exception:
                pass

        if updates.entry_type is not None:
            entry.entry_type = updates.entry_type

        if updates.metadata is not None:
            entry.metadata_json = updates.metadata

        if updates.sharing_scope is not None:
            entry.sharing_scope = updates.sharing_scope.value if hasattr(updates.sharing_scope, 'value') else updates.sharing_scope

        await self.db.flush()
        await self.db.refresh(entry)

        return self._entry_to_response(entry)

    async def get_graph_workbench(
        self,
        tenant_id: UUID,
        project_id: UUID | None = None,
        source_file_id: UUID | None = None,
        entry_type: str | None = None,
        limit: int = 200,
        user_id: UUID | None = None,
    ) -> KnowledgeGraphWorkbenchResponse:
        """Return graph nodes, edges, and readiness metrics for the UI."""
        query = select(KnowledgeEntry).where(
            KnowledgeEntry.tenant_id == tenant_id,
            KnowledgeEntry.deleted_at.is_(None),
        )

        if user_id:
            query = query.where(
                or_(
                    KnowledgeEntry.sharing_scope.in_(["project", "tenant", "global"]),
                    and_(
                        KnowledgeEntry.sharing_scope == "private",
                        KnowledgeEntry.created_by_id == user_id,
                    ),
                )
            )
        if project_id:
            query = query.where(KnowledgeEntry.project_id == project_id)
        if source_file_id:
            query = query.where(KnowledgeEntry.source_file_id == source_file_id)
        if entry_type:
            query = query.where(KnowledgeEntry.entry_type == entry_type)

        query = query.order_by(KnowledgeEntry.created_at.desc()).limit(limit)
        result = await self.db.execute(query)
        entries = list(result.scalars().all())
        entry_ids = [entry.id for entry in entries]

        links: list[KnowledgeLink] = []
        if entry_ids:
            link_query = (
                select(KnowledgeLink)
                .where(
                    KnowledgeLink.tenant_id == tenant_id,
                    KnowledgeLink.deleted_at.is_(None),
                    KnowledgeLink.source_entry_id.in_(entry_ids),
                    KnowledgeLink.target_entry_id.in_(entry_ids),
                )
                .order_by(KnowledgeLink.created_at.desc())
            )
            link_result = await self.db.execute(link_query)
            links = list(link_result.scalars().all())

        nodes = [self._entry_to_graph_node(entry) for entry in entries]
        edges = [self._link_to_graph_edge(link) for link in links]
        linked_ids = {edge.source for edge in edges} | {edge.target for edge in edges}
        gaps = [node for node in nodes if node.status in {"gap", "conflict"}]
        isolated_nodes = [node for node in nodes if node.id not in linked_ids]

        summary = KnowledgeGraphSummary(
            node_count=len(nodes),
            edge_count=len(edges),
            source_count=len({node.source_file_id or node.source_label for node in nodes}),
            ready_count=len([node for node in nodes if node.status == "ready"]),
            gap_count=len([node for node in nodes if node.status == "gap"]),
            conflict_count=len([node for node in nodes if node.status == "conflict"]),
            isolated_count=len(isolated_nodes),
        )

        return KnowledgeGraphWorkbenchResponse(
            nodes=nodes,
            edges=edges,
            summary=summary,
            gaps=gaps,
            isolated_nodes=isolated_nodes,
        )

    async def delete_entry(self, entry_id: UUID, tenant_id: UUID) -> bool:
        """Soft delete a knowledge entry.

        Args:
            entry_id: Entry UUID
            tenant_id: Tenant UUID for verification

        Returns:
            True if deleted, False if not found
        """
        result = await self.db.execute(
            select(KnowledgeEntry).where(
                KnowledgeEntry.id == entry_id,
                KnowledgeEntry.tenant_id == tenant_id,
                KnowledgeEntry.deleted_at.is_(None),
            )
        )
        entry = result.scalar_one_or_none()

        if not entry:
            return False

        entry.deleted_at = datetime.now(timezone.utc)

        # Soft delete from FTS
        try:
            await self.search_provider.delete_document(entry_id)
        except Exception:
            pass

        # Delete from vector store
        try:
            await self.vector_provider.delete_vector(entry_id)
        except Exception:
            pass

        await self.db.flush()
        return True

    async def link_entries(
        self,
        source_id: UUID,
        target_id: UUID,
        link_type: str,
        tenant_id: UUID,
        confidence: float | None = 1.0,
        metadata: dict[str, Any] | None = None,
    ) -> KnowledgeLinkResponse | None:
        """Create a knowledge link between two entries.

        Args:
            source_id: Source entry UUID
            target_id: Target entry UUID
            link_type: Type of link (cites, extends, depends_on, implements)
            tenant_id: Tenant UUID for verification
            confidence: Optional confidence score
            metadata: Optional metadata

        Returns:
            KnowledgeLinkResponse if created, None if entries not found
        """
        # Verify both entries exist and belong to tenant
        source_result = await self.db.execute(
            select(KnowledgeEntry).where(
                KnowledgeEntry.id == source_id,
                KnowledgeEntry.tenant_id == tenant_id,
                KnowledgeEntry.deleted_at.is_(None),
            )
        )
        source = source_result.scalar_one_or_none()

        if not source:
            return None

        target_result = await self.db.execute(
            select(KnowledgeEntry).where(
                KnowledgeEntry.id == target_id,
                KnowledgeEntry.tenant_id == tenant_id,
                KnowledgeEntry.deleted_at.is_(None),
            )
        )
        target = target_result.scalar_one_or_none()

        if not target:
            return None

        # Create link via graph provider
        link = await self.graph_provider.add_edge(
            source_id,
            target_id,
            link_type,
            metadata,
        )

        if link:
            link.confidence = confidence

        await self.db.flush()

        return KnowledgeLinkResponse(
            id=link.id,
            tenant_id=link.tenant_id,
            source_entry_id=link.source_entry_id,
            target_entry_id=link.target_entry_id,
            link_type=link.link_type,
            confidence=link.confidence,
            metadata_json=link.metadata_json,
            created_at=link.created_at,
        )

    async def get_entry_links(
        self,
        entry_id: UUID,
        tenant_id: UUID,
        direction: str = "both",
    ) -> list[KnowledgeLinkResponse]:
        """Get all links for an entry.

        Args:
            entry_id: Entry UUID
            tenant_id: Tenant UUID for verification
            direction: Link direction ("outgoing", "incoming", "both")

        Returns:
            List of KnowledgeLinkResponse
        """
        # Verify entry exists
        entry_result = await self.db.execute(
            select(KnowledgeEntry).where(
                KnowledgeEntry.id == entry_id,
                KnowledgeEntry.tenant_id == tenant_id,
                KnowledgeEntry.deleted_at.is_(None),
            )
        )
        if not entry_result.scalar_one_or_none():
            return []

        links = []
        query = select(KnowledgeLink).where(KnowledgeLink.tenant_id == tenant_id)

        if direction == "outgoing":
            query = query.where(KnowledgeLink.source_entry_id == entry_id)
        elif direction == "incoming":
            query = query.where(KnowledgeLink.target_entry_id == entry_id)
        else:
            query = query.where(
                or_(
                    KnowledgeLink.source_entry_id == entry_id,
                    KnowledgeLink.target_entry_id == entry_id,
                )
            )

        result = await self.db.execute(query)
        for link in result.scalars().all():
            links.append(KnowledgeLinkResponse(
                id=link.id,
                tenant_id=link.tenant_id,
                source_entry_id=link.source_entry_id,
                target_entry_id=link.target_entry_id,
                link_type=link.link_type,
                confidence=link.confidence,
                metadata_json=link.metadata_json,
                created_at=link.created_at,
            ))

        return links

    async def update_link(
        self,
        link_id: UUID,
        tenant_id: UUID,
        link_type: str | None = None,
        confidence: float | None = None,
        metadata: dict[str, Any] | None = None,
    ) -> KnowledgeLinkResponse | None:
        """Update a knowledge graph relationship."""
        result = await self.db.execute(
            select(KnowledgeLink).where(
                KnowledgeLink.id == link_id,
                KnowledgeLink.tenant_id == tenant_id,
                KnowledgeLink.deleted_at.is_(None),
            )
        )
        link = result.scalar_one_or_none()
        if not link:
            return None

        if link_type is not None:
            link.link_type = link_type
        if confidence is not None:
            link.confidence = confidence
        if metadata is not None:
            link.metadata_json = metadata

        await self.db.flush()
        await self.db.refresh(link)

        return KnowledgeLinkResponse(
            id=link.id,
            tenant_id=link.tenant_id,
            source_entry_id=link.source_entry_id,
            target_entry_id=link.target_entry_id,
            link_type=link.link_type,
            confidence=link.confidence,
            metadata_json=link.metadata_json,
            created_at=link.created_at,
        )

    async def delete_link(self, link_id: UUID, tenant_id: UUID) -> bool:
        """Soft delete a knowledge graph relationship."""
        result = await self.db.execute(
            select(KnowledgeLink).where(
                KnowledgeLink.id == link_id,
                KnowledgeLink.tenant_id == tenant_id,
                KnowledgeLink.deleted_at.is_(None),
            )
        )
        link = result.scalar_one_or_none()
        if not link:
            return False

        link.deleted_at = datetime.now(timezone.utc)
        await self.db.flush()
        return True

    async def get_lineage(
        self,
        entry_id: UUID,
        tenant_id: UUID,
    ) -> list[LineageRecordResponse]:
        """Get lineage records for an entry.

        Args:
            entry_id: Entry UUID
            tenant_id: Tenant UUID for verification

        Returns:
            List of LineageRecordResponse
        """
        result = await self.db.execute(
            select(LineageRecord).where(
                LineageRecord.tenant_id == tenant_id,
                or_(
                    and_(LineageRecord.source_type == "knowledge_entry", LineageRecord.source_id == entry_id),
                    and_(LineageRecord.target_type == "knowledge_entry", LineageRecord.target_id == entry_id),
                ),
            ).order_by(LineageRecord.created_at.desc())
        )

        return [
            LineageRecordResponse(
                id=lr.id,
                tenant_id=lr.tenant_id,
                project_id=lr.project_id,
                source_type=lr.source_type,
                source_id=lr.source_id,
                target_type=lr.target_type,
                target_id=lr.target_id,
                lineage_type=lr.lineage_type,
                metadata_json=lr.metadata_json,
                created_at=lr.created_at,
            )
            for lr in result.scalars().all()
        ]

    async def add_lineage(
        self,
        tenant_id: UUID,
        project_id: UUID,
        source_type: str,
        source_id: UUID,
        target_type: str,
        target_id: UUID,
        lineage_type: str,
        metadata: dict[str, Any] | None = None,
    ) -> LineageRecordResponse:
        """Add a lineage record.

        Args:
            tenant_id: Tenant UUID
            project_id: Project UUID
            source_type: Source entity type
            source_id: Source entity UUID
            target_type: Target entity type
            target_id: Target entity UUID
            lineage_type: Lineage relationship type
            metadata: Optional metadata

        Returns:
            Created LineageRecordResponse
        """
        lineage = LineageRecord(
            tenant_id=tenant_id,
            project_id=project_id,
            source_type=source_type,
            source_id=source_id,
            target_type=target_type,
            target_id=target_id,
            lineage_type=lineage_type,
            metadata_json=metadata,
        )
        self.db.add(lineage)
        await self.db.flush()
        await self.db.refresh(lineage)

        return LineageRecordResponse(
            id=lineage.id,
            tenant_id=lineage.tenant_id,
            project_id=lineage.project_id,
            source_type=lineage.source_type,
            source_id=lineage.source_id,
            target_type=lineage.target_type,
            target_id=lineage.target_id,
            lineage_type=lineage.lineage_type,
            metadata_json=lineage.metadata_json,
            created_at=lineage.created_at,
        )

    async def get_provenance(
        self,
        entry_id: UUID,
        tenant_id: UUID,
    ) -> list[ProvenanceRecordResponse]:
        """Get provenance records for an entry.

        Args:
            entry_id: Entry UUID
            tenant_id: Tenant UUID for verification

        Returns:
            List of ProvenanceRecordResponse
        """
        result = await self.db.execute(
            select(ProvenanceRecord).where(
                ProvenanceRecord.entry_id == entry_id,
                ProvenanceRecord.tenant_id == tenant_id,
            ).order_by(ProvenanceRecord.created_at.desc())
        )

        return [
            ProvenanceRecordResponse(
                id=pr.id,
                tenant_id=pr.tenant_id,
                project_id=pr.project_id,
                entry_id=pr.entry_id,
                provider_id=pr.provider_id,
                provider_version_id=pr.provider_version_id,
                raw_artifact_id=pr.raw_artifact_id,
                confidence=pr.confidence,
                normalization_notes=pr.normalization_notes,
                created_at=pr.created_at,
            )
            for pr in result.scalars().all()
        ]

    async def add_provenance(
        self,
        tenant_id: UUID,
        project_id: UUID,
        entry_id: UUID,
        provider_id: str,
        provider_version_id: str | None = None,
        raw_artifact_id: str | None = None,
        confidence: float | None = 1.0,
        normalization_notes: str | None = None,
    ) -> ProvenanceRecordResponse:
        """Add a provenance record.

        Args:
            tenant_id: Tenant UUID
            project_id: Project UUID
            entry_id: Entry UUID
            provider_id: Provider identifier
            provider_version_id: Optional provider version
            raw_artifact_id: Optional raw artifact ID
            confidence: Optional confidence score
            normalization_notes: Optional notes

        Returns:
            Created ProvenanceRecordResponse
        """
        provenance = ProvenanceRecord(
            tenant_id=tenant_id,
            project_id=project_id,
            entry_id=entry_id,
            provider_id=provider_id,
            provider_version_id=provider_version_id,
            raw_artifact_id=raw_artifact_id,
            confidence=confidence,
            normalization_notes=normalization_notes,
        )
        self.db.add(provenance)
        await self.db.flush()
        await self.db.refresh(provenance)

        return ProvenanceRecordResponse(
            id=provenance.id,
            tenant_id=provenance.tenant_id,
            project_id=provenance.project_id,
            entry_id=provenance.entry_id,
            provider_id=provenance.provider_id,
            provider_version_id=provenance.provider_version_id,
            raw_artifact_id=provenance.raw_artifact_id,
            confidence=provenance.confidence,
            normalization_notes=provenance.normalization_notes,
            created_at=provenance.created_at,
        )

    async def ingest_source_file(
        self,
        source_file_id: UUID,
        tenant_id: UUID,
        project_id: UUID,
        storage: Any | None = None,
    ) -> list[KnowledgeEntryResponse]:
        """Ingest a source file into knowledge entries.

        Processes a source file and creates knowledge entries from its content.
        Supports markdown, docx, pptx, and txt files.

        Args:
            source_file_id: Source file UUID
            tenant_id: Tenant UUID
            project_id: Project UUID

        Returns:
            List of created KnowledgeEntryResponse
        """
        from app.domains.projects.models import SourceFile

        result = await self.db.execute(
            select(SourceFile).where(
                SourceFile.id == source_file_id,
                SourceFile.tenant_id == tenant_id,
                SourceFile.deleted_at.is_(None),
            )
        )
        source_file = result.scalar_one_or_none()

        if not source_file:
            return []

        # Read file content from storage
        storage = storage or get_storage_provider()
        storage_handle = StorageHandle(
            path=source_file.storage_path,
            filename=source_file.original_filename,
            content_type=source_file.content_type,
            size=int(source_file.size),
            hash=source_file.hash,
            storage_backend=settings.STORAGE_BACKEND,
        )

        content_bytes = await storage.download(storage_handle)
        content = self._extract_text_from_bytes(
            content_bytes,
            source_file.content_type,
            source_file.original_filename,
        )

        chunks = self._source_ingestion_chunks(content)

        entries = []
        for i, chunk in enumerate(chunks):
            chunk_content = chunk["content"]
            entry = await self.create_entry(
                tenant_id=tenant_id,
                project_id=project_id,
                entry_type="text",
                content=chunk_content,
                source_file_id=source_file_id,
                metadata={
                    "filename": source_file.original_filename,
                    "content_type": source_file.content_type,
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "title": chunk["title"],
                    "summary": chunk_content[:240],
                    "ingestion": {
                        "source_file_status": "ready",
                        "source_file_id": str(source_file_id),
                        "source_file_name": source_file.original_filename,
                        "chunk_index": i,
                        "chunk_count": len(chunks),
                        "identifier": chunk["identifier"],
                        "section_type": chunk["section_type"],
                        "chunking_strategy": chunk["chunking_strategy"],
                    },
                },
            )
            entries.append(entry)

            await self.add_provenance(
                tenant_id=tenant_id,
                project_id=project_id,
                entry_id=entry.id,
                provider_id="source_file_ingest",
                provider_version_id="1.0",
                raw_artifact_id=str(source_file_id),
                normalization_notes=f"chunk {i + 1}/{len(chunks)} via {chunk['chunking_strategy']}",
            )

        for i in range(1, len(entries)):
            await self.link_entries(
                source_id=entries[i].id,
                target_id=entries[i - 1].id,
                link_type="depends_on",
                tenant_id=tenant_id,
                confidence=0.8,
                metadata={
                    "ingestion": {
                        "source_file_id": str(source_file_id),
                        "source_file_name": source_file.original_filename,
                        "source_chunk_index": i,
                        "target_chunk_index": i - 1,
                        "rule": "structured_sequence",
                    }
                },
            )

        return entries

    def _source_ingestion_chunks(self, text: str) -> list[dict[str, Any]]:
        structured_chunks = self._structured_source_chunks(text)
        if structured_chunks:
            return structured_chunks

        raw_chunks = self._chunk_text(text, max_chunk_size=4000)
        return [
            {
                "identifier": None,
                "section_type": None,
                "title": f"Chunk {index + 1}",
                "content": chunk,
                "chunking_strategy": "size",
            }
            for index, chunk in enumerate(raw_chunks)
        ]

    def _structured_source_chunks(self, text: str) -> list[dict[str, Any]]:
        pattern = re.compile(r"^\s*([A-Z]{2,12}-\d{1,8})\s*[:：-]\s*(.+?)\s*$")
        chunks: list[dict[str, Any]] = []
        current: dict[str, Any] | None = None

        for raw_line in text.splitlines():
            line = raw_line.strip()
            if not line:
                continue

            match = pattern.match(line)
            if match:
                if current:
                    chunks.append(current)
                identifier = match.group(1)
                section_type = identifier.split("-", 1)[0].lower()
                current = {
                    "identifier": identifier,
                    "section_type": section_type,
                    "title": identifier,
                    "content": line,
                    "chunking_strategy": "structured_identifier",
                }
                continue

            if current:
                current["content"] = f"{current['content']}\n{line}"

        if current:
            chunks.append(current)

        return chunks if len(chunks) >= 2 else []

    def _extract_text_from_bytes(
        self,
        content: bytes,
        content_type: str,
        filename: str,
    ) -> str:
        """Extract text content from file bytes based on content type.

        Args:
            content: File content as bytes
            content_type: MIME type
            filename: Original filename

        Returns:
            Extracted text content
        """
        ext = filename.lower().split('.')[-1] if '.' in filename else ''

        # Handle text/markdown directly
        if content_type in ('text/markdown', 'text/plain') or ext in ('md', 'txt', 'markdown'):
            try:
                return content.decode('utf-8')
            except UnicodeDecodeError:
                return content.decode('utf-8', errors='replace')

        # Handle DOCX files
        if content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or ext == 'docx':
            try:
                import io
                from docx import Document
                doc = Document(io.BytesIO(content))
                extracted = '\n\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
                if not extracted.strip():
                    raise SourceTextExtractionError(f"No extractable DOCX text: {filename}")
                return extracted
            except ImportError as exc:
                raise SourceTextExtractionError("DOCX parser dependency is not available") from exc
            except SourceTextExtractionError:
                raise
            except Exception as exc:
                raise SourceTextExtractionError(f"Failed to parse DOCX: {filename}") from exc

        # Handle PPTX files
        if content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or ext == 'pptx':
            try:
                import io
                from pptx import Presentation
                prs = Presentation(io.BytesIO(content))
                text_content = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            text_content.append(shape.text)
                extracted = '\n\n'.join(text_content)
                if not extracted.strip():
                    raise SourceTextExtractionError(f"No extractable PPTX text: {filename}")
                return extracted
            except ImportError as exc:
                raise SourceTextExtractionError("PPTX parser dependency is not available") from exc
            except SourceTextExtractionError:
                raise
            except Exception as exc:
                raise SourceTextExtractionError(f"Failed to parse PPTX: {filename}") from exc

        # Handle XLSX files
        if content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' or ext == 'xlsx':
            try:
                import io
                from openpyxl import load_workbook

                workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
                rows: list[str] = []
                for worksheet in workbook.worksheets:
                    rows.append(f"# {worksheet.title}")
                    for row in worksheet.iter_rows(values_only=True):
                        values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                        if values:
                            rows.append(" | ".join(values))
                extracted = "\n".join(rows)
                if not extracted.strip():
                    raise SourceTextExtractionError(f"No extractable XLSX text: {filename}")
                return extracted
            except ImportError as exc:
                raise SourceTextExtractionError("XLSX parser dependency is not available") from exc
            except SourceTextExtractionError:
                raise
            except Exception as exc:
                raise SourceTextExtractionError(f"Failed to parse XLSX: {filename}") from exc

        if content_type == 'application/pdf' or ext == 'pdf':
            try:
                import io
                from pypdf import PdfReader

                reader = PdfReader(io.BytesIO(content))
                extracted = "\n\n".join(
                    page_text.strip()
                    for page in reader.pages
                    if (page_text := page.extract_text() or "").strip()
                )
                if not extracted.strip():
                    raise SourceTextExtractionError(f"No extractable PDF text: {filename}")
                return extracted
            except ImportError as exc:
                raise SourceTextExtractionError("PDF parser dependency is not available") from exc
            except SourceTextExtractionError:
                raise
            except Exception as exc:
                raise SourceTextExtractionError(f"Failed to parse PDF: {filename}") from exc

        # Fallback for unknown text-compatible types
        try:
            extracted = content.decode('utf-8')
            if not extracted.strip():
                raise SourceTextExtractionError(f"No extractable text: {filename}")
            return extracted
        except UnicodeDecodeError as exc:
            raise SourceTextExtractionError(f"Unsupported binary source content: {filename}") from exc

    def _chunk_text(self, text: str, max_chunk_size: int = 4000) -> list[str]:
        """Split text into chunks of maximum size.

        Args:
            text: Text to chunk
            max_chunk_size: Maximum characters per chunk

        Returns:
            List of text chunks
        """
        if len(text) <= max_chunk_size:
            return [text]

        chunks = []
        paragraphs = text.split('\n\n')
        current_chunk = []
        current_size = 0

        for para in paragraphs:
            para_size = len(para) + 2  # +2 for \n\n
            if current_size + para_size > max_chunk_size and current_chunk:
                chunks.append('\n\n'.join(current_chunk))
                current_chunk = []
                current_size = 0
            current_chunk.append(para)
            current_size += para_size

        if current_chunk:
            chunks.append('\n\n'.join(current_chunk))

        return chunks

    async def bulk_create_entries(
        self,
        tenant_id: UUID,
        project_id: UUID,
        entries: list[KnowledgeEntryCreate],
        generate_embeddings: bool = True,
    ) -> BulkEntryResponse:
        """Bulk create knowledge entries.

        Args:
            tenant_id: Tenant UUID
            project_id: Project UUID
            entries: List of entry creation schemas
            generate_embeddings: Whether to generate embeddings

        Returns:
            BulkEntryResponse with counts
        """
        async def create_one(entry_data: KnowledgeEntryCreate) -> tuple[bool, str | None]:
            """Create a single entry and return success status and error message."""
            try:
                await self.create_entry(
                    tenant_id=tenant_id,
                    project_id=project_id,
                    entry_type=entry_data.entry_type,
                    content=entry_data.content,
                    source_file_id=entry_data.source_file_id,
                    metadata=entry_data.metadata,
                    generate_embedding=entry_data.generate_embedding and generate_embeddings,
                )
                return (True, None)
            except Exception as e:
                return (False, str(e))

        # Process all entries concurrently
        results = await asyncio.gather(*[create_one(e) for e in entries])

        created = sum(1 for r in results if r[0])
        failed = sum(1 for r in results if not r[0])
        errors = [r[1] for r in results if r[1]]

        return BulkEntryResponse(
            created=created,
            failed=failed,
            errors=errors,
        )

    async def get_neighbors(
        self,
        entry_id: UUID,
        tenant_id: UUID,
        depth: int = 1,
        direction: str = "both",
    ) -> list[GraphNeighborResponse]:
        """Get neighboring entries via graph traversal.

        Args:
            entry_id: Entry UUID
            tenant_id: Tenant UUID for verification
            depth: Traversal depth
            direction: Traversal direction

        Returns:
            List of GraphNeighborResponse
        """
        neighbors = await self.graph_provider.get_neighbors(entry_id, depth, direction)

        results = []
        for entry, link, d in neighbors:
            if entry.tenant_id == tenant_id:
                response = KnowledgeEntryResponse(
                    id=entry.id,
                    tenant_id=entry.tenant_id,
                    project_id=entry.project_id,
                    source_file_id=entry.source_file_id,
                    entry_type=entry.entry_type,
                    content=entry.content,
                    content_hash=entry.content_hash,
                    vector_embedding=entry.vector_embedding,
                    metadata=entry.metadata_json,
                    created_at=entry.created_at,
                    updated_at=entry.updated_at,
                    deleted_at=entry.deleted_at,
                )
                link_type = link.link_type if link else "unknown"
                results.append(GraphNeighborResponse(
                    entry=response,
                    link_type=link_type,
                    depth=d,
                ))

        return results
