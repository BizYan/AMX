"""Export Domain Service

Business logic for document export operations to Word, Markdown, and PPTX formats.
"""

import hashlib
import io
import re
from datetime import datetime, timezone
from typing import Any
from uuid import UUID

from sqlalchemy import select, func
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import selectinload

from app.domains.export.models import ExportJob, ExportArtifact, ExportStatus, ExportType
from app.domains.export.schemas import (
    WordExportRequest,
    MarkdownExportRequest,
    PPTXExportRequest,
    ExportReadinessDocumentBlocker,
    ExportReadinessRequiredType,
    ExportReadinessResponse,
    ExportReleaseEvidenceResponse,
    ExportReleaseEvidenceSummary,
    ExportReleaseGate,
    ExportReleasePriorityAction,
    ExportReleaseRiskItem,
    ExportJobWithArtifactsResponse,
    ExportArtifactResponse,
)
from app.core.placeholders import (
    PLACEHOLDER_PATTERN as SHARED_PLACEHOLDER_PATTERN,
    contains_placeholder,
    extract_placeholders,
    substitute_placeholders,
)
from app.services.storage import StorageHandle, get_storage_provider


class ExportService:
    """Service for document export operations.

    Handles export job management, actual document export to various formats,
    and artifact storage.
    """

    def __init__(self, db: AsyncSession):
        """Initialize export service.

        Args:
            db: Async database session
        """
        self.db = db

    async def get_project_export_readiness(
        self,
        project_id: UUID,
        tenant_id: UUID,
    ) -> ExportReadinessResponse:
        """Return production delivery package readiness for a project."""
        from app.domains.documents.models import Document, DocumentStatus, DocumentType
        from app.models.projects import Project

        project_result = await self.db.execute(
            select(Project.id).where(
                Project.id == project_id,
                Project.tenant_id == tenant_id,
                Project.deleted_at.is_(None),
            )
        )
        if project_result.scalar_one_or_none() is None:
            raise ValueError(f"Project not found: {project_id}")

        result = await self.db.execute(
            select(Document).where(
                Document.project_id == project_id,
                Document.tenant_id == tenant_id,
                Document.deleted_at.is_(None),
            )
        )
        documents = list(result.scalars().all())
        production_statuses = {
            DocumentStatus.APPROVED.value,
            DocumentStatus.PUBLISHED.value,
        }
        required_types = [
            (DocumentType.URS.value, "URS"),
            (DocumentType.BRD.value, "BRD"),
            (DocumentType.PRD.value, "PRD"),
            (DocumentType.DETAILED_DESIGN.value, "设计说明"),
            (DocumentType.TEST_CASE.value, "测试用例"),
        ]

        blockers: list[ExportReadinessDocumentBlocker] = []
        ready_by_type: dict[str, int] = {}
        blocked_by_type: dict[str, int] = {}
        blocked_document_ids: set[UUID] = set()

        for document in documents:
            metadata = document.metadata_json or {}
            generation_status = str(metadata.get("generation_status") or "").lower()
            has_placeholder = (
                generation_status in {"placeholder", "partial", "failed"}
                or metadata.get("has_placeholders") is True
                or document.status == "placeholder"
            )
            delivery = metadata.get("delivery") if isinstance(metadata, dict) else None
            delivery_readiness = delivery.get("delivery_readiness") if isinstance(delivery, dict) else None
            has_delivery_readiness_failure = (
                isinstance(delivery_readiness, dict)
                and delivery_readiness.get("ready") is False
            )
            empty_content = not (document.content or "").strip()
            not_formal = document.status not in production_statuses
            document_blockers: list[tuple[str, str]] = []

            if has_placeholder:
                document_blockers.append(
                    ("文档仍包含占位内容", "回到项目文档工作台补齐内容或重新生成后再发布。")
                )
            if has_delivery_readiness_failure:
                document_blockers.append(
                    (
                        "Document delivery readiness failed",
                        "Resolve document delivery readiness blockers before export.",
                    )
                )
            unresolved_comments = await self._count_unresolved_document_comments(document.id, tenant_id)
            if unresolved_comments:
                document_blockers.append(
                    (
                        f"Document has {unresolved_comments} unresolved comments",
                        "Resolve all review comments before exporting the delivery package.",
                    )
                )
            failed_quality_checks = await self._count_failed_document_quality_checks(document.id, tenant_id)
            if failed_quality_checks:
                document_blockers.append(
                    (
                        f"Document quality check failed ({failed_quality_checks})",
                        "Rerun quality review and resolve failed quality findings before export.",
                    )
                )
            if empty_content:
                document_blockers.append(
                    ("文档内容为空", "补齐正文并生成版本快照后再纳入交付包。")
                )
            if not_formal:
                document_blockers.append(
                    ("文档尚未批准或发布", "完成评审、批准或发布流程后再用于正式交付。")
                )

            if document_blockers:
                blocked_document_ids.add(document.id)
                blocked_by_type[document.doc_type] = blocked_by_type.get(document.doc_type, 0) + 1
                for reason, recommended_action in document_blockers:
                    blockers.append(
                        ExportReadinessDocumentBlocker(
                            document_id=document.id,
                            title=document.title,
                            doc_type=document.doc_type,
                            status=document.status,
                            reason=reason,
                            recommended_action=recommended_action,
                        )
                    )
            else:
                ready_by_type[document.doc_type] = ready_by_type.get(document.doc_type, 0) + 1

        required_type_states: list[ExportReadinessRequiredType] = []
        missing_required_types: list[str] = []
        for doc_type, label in required_types:
            ready_count = ready_by_type.get(doc_type, 0)
            blocked_count = blocked_by_type.get(doc_type, 0)
            if ready_count > 0:
                status = "ready"
            elif blocked_count > 0:
                status = "blocked"
            else:
                status = "missing"
                missing_required_types.append(doc_type)
            required_type_states.append(
                ExportReadinessRequiredType(
                    doc_type=doc_type,
                    label=label,
                    ready_count=ready_count,
                    blocked_count=blocked_count,
                    status=status,
                )
            )

        ready_required_count = sum(1 for item in required_type_states if item.status == "ready")
        readiness_score = round((ready_required_count / len(required_type_states)) * 100)
        if blockers:
            readiness_score = max(0, readiness_score - min(20, len(blockers) * 4))

        recommended_actions: list[str] = []
        if missing_required_types:
            missing_labels = [
                item.label for item in required_type_states if item.status == "missing"
            ]
            recommended_actions.append(
                f"补齐核心交付文档类型：{', '.join(missing_labels)}。"
            )
        if blockers:
            recommended_actions.append("处理占位、空内容和未发布文档后再生成正式交付包。")
        if not documents:
            recommended_actions.append("先在项目文档中生成并发布核心交付文档。")
        if not recommended_actions:
            recommended_actions.append("可以生成正式交付包，并保留审计清单和下载产物。")

        exportable_documents = sum(ready_by_type.values())
        return ExportReadinessResponse(
            project_id=project_id,
            total_documents=len(documents),
            exportable_documents=exportable_documents,
            blocked_documents=len(blocked_document_ids),
            readiness_score=readiness_score,
            can_export_production=(
                exportable_documents > 0 and not missing_required_types and not blockers
            ),
            missing_required_types=missing_required_types,
            required_types=required_type_states,
            blockers=blockers,
            recommended_actions=recommended_actions,
        )

    async def get_project_release_evidence(
        self,
        project_id: UUID,
        tenant_id: UUID,
    ) -> ExportReleaseEvidenceResponse:
        """Aggregate project package readiness, export jobs, artifacts, and release actions."""
        readiness = await self.get_project_export_readiness(
            project_id=project_id,
            tenant_id=tenant_id,
        )
        result = await self.db.execute(
            select(ExportJob)
            .options(selectinload(ExportJob.artifacts))
            .where(
                ExportJob.project_id == project_id,
                ExportJob.tenant_id == tenant_id,
            )
            .order_by(ExportJob.created_at.desc())
        )
        jobs = list(result.scalars().unique().all())
        completed_jobs = [job for job in jobs if job.status == ExportStatus.COMPLETED.value]
        failed_jobs = [job for job in jobs if job.status == ExportStatus.FAILED.value]
        package_jobs = [job for job in completed_jobs if job.export_type == ExportType.PROJECT_PACKAGE.value]
        latest_job = package_jobs[0] if package_jobs else (completed_jobs[0] if completed_jobs else None)
        artifacts = [
            artifact
            for job in completed_jobs
            for artifact in (job.artifacts or [])
        ]
        recent_artifacts = sorted(
            artifacts,
            key=lambda artifact: (artifact.created_at, artifact.filename),
            reverse=True,
        )[:6]
        covered_formats = sorted({self._artifact_format(artifact.filename) for artifact in artifacts})
        covered_formats = [item for item in covered_formats if item in {"Markdown", "PPTX", "Word"}]
        missing_formats = [
            format_name
            for format_name in ["Markdown", "PPTX", "Word"]
            if format_name not in covered_formats
        ]
        latest_completed_at = max(
            (job.completed_at for job in completed_jobs if job.completed_at),
            default=None,
        )
        summary = ExportReleaseEvidenceSummary(
            total_jobs=len(jobs),
            completed_jobs=len(completed_jobs),
            failed_jobs=len(failed_jobs),
            artifact_count=len(artifacts),
            covered_formats=covered_formats,
            missing_formats=missing_formats,
            production_package_jobs=len(package_jobs),
            latest_completed_at=latest_completed_at,
        )
        risk_items = self._build_release_risks(readiness, summary)
        release_gate = self._build_release_gate(readiness, summary, risk_items)
        priority_actions = self._build_release_actions(readiness, summary)

        return ExportReleaseEvidenceResponse(
            project_id=project_id,
            release_gate=release_gate,
            readiness=readiness,
            summary=summary,
            latest_job=ExportJobWithArtifactsResponse.model_validate(latest_job) if latest_job else None,
            recent_artifacts=[ExportArtifactResponse.model_validate(artifact) for artifact in recent_artifacts],
            risk_items=risk_items,
            priority_actions=priority_actions,
        )

    def _build_release_risks(
        self,
        readiness: ExportReadinessResponse,
        summary: ExportReleaseEvidenceSummary,
    ) -> list[ExportReleaseRiskItem]:
        risks: list[ExportReleaseRiskItem] = []
        if not readiness.can_export_production:
            risks.append(ExportReleaseRiskItem(
                code="export_readiness_blocked",
                severity="critical",
                title="核心交付清单未就绪",
                detail="仍有缺失、占位、空内容或未批准发布的核心文档，不能作为正式交付包发布。",
                count=len(readiness.missing_required_types) + readiness.blocked_documents,
                href="/documents",
            ))
        if summary.failed_jobs:
            risks.append(ExportReleaseRiskItem(
                code="failed_exports",
                severity="high",
                title="存在失败导出任务",
                detail="导出失败需要确认变量、模板或产物生成错误，并重新生成交付包。",
                count=summary.failed_jobs,
                href="/exports",
            ))
        if summary.missing_formats:
            risks.append(ExportReleaseRiskItem(
                code="missing_format_coverage",
                severity="medium",
                title="交付格式覆盖不完整",
                detail=f"尚未生成这些正式交付格式：{', '.join(summary.missing_formats)}。",
                count=len(summary.missing_formats),
                href="/exports",
            ))
        if summary.production_package_jobs == 0:
            risks.append(ExportReleaseRiskItem(
                code="missing_project_package",
                severity="high",
                title="尚无正式项目交付包",
                detail="需要至少生成一次项目级交付包并保留下载产物。",
                count=1,
                href="/exports",
            ))
        return risks

    def _build_release_gate(
        self,
        readiness: ExportReadinessResponse,
        summary: ExportReleaseEvidenceSummary,
        risk_items: list[ExportReleaseRiskItem],
    ) -> ExportReleaseGate:
        blockers = [item.title for item in risk_items if item.severity in {"critical", "high"}]
        warnings = [item.title for item in risk_items if item.severity == "medium"]
        if blockers:
            return ExportReleaseGate(
                status="blocked",
                label="发布阻断",
                summary="导出发布证据仍存在必须处理的阻断项。",
                blockers=blockers,
                warnings=warnings,
            )
        if warnings or not readiness.can_export_production or summary.missing_formats:
            return ExportReleaseGate(
                status="attention",
                label="需复核",
                summary="核心阻断已清除，但正式发布前仍需补齐证据或格式覆盖。",
                blockers=[],
                warnings=warnings,
            )
        return ExportReleaseGate(
            status="passed",
            label="可进入发布",
            summary="核心交付文档、导出任务和交付格式证据已满足发布要求。",
            blockers=[],
            warnings=[],
        )

    def _build_release_actions(
        self,
        readiness: ExportReadinessResponse,
        summary: ExportReleaseEvidenceSummary,
    ) -> list[ExportReleasePriorityAction]:
        actions: list[ExportReleasePriorityAction] = []
        if not readiness.can_export_production:
            actions.append(ExportReleasePriorityAction(
                code="fix_delivery_documents",
                title="补齐核心交付文档",
                description="处理缺失文档、占位内容、空正文和未批准发布状态后再生成正式交付包。",
                href="/documents",
                priority="critical",
            ))
        if summary.failed_jobs:
            actions.append(ExportReleasePriorityAction(
                code="triage_failed_exports",
                title="处置失败导出任务",
                description="查看失败原因，补齐模板变量或格式配置后重新导出。",
                href="/exports",
                priority="high",
            ))
        if summary.missing_formats:
            actions.append(ExportReleasePriorityAction(
                code="generate_missing_formats",
                title="补齐交付格式",
                description=f"生成缺失格式：{', '.join(summary.missing_formats)}。",
                href="/exports",
                priority="medium",
            ))
        if not actions:
            actions.append(ExportReleasePriorityAction(
                code="preserve_release_evidence",
                title="保留发布证据",
                description="下载交付包、审计清单和任务记录，作为发布验收附件。",
                href="/exports",
                priority="medium",
            ))
        return actions

    def _artifact_format(self, filename: str) -> str:
        extension = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if extension == "docx":
            return "Word"
        if extension == "md":
            return "Markdown"
        if extension == "pptx":
            return "PPTX"
        return extension.upper() if extension else "File"

    def _document_delivery_readiness_failed(self, document: Any) -> bool:
        metadata = getattr(document, "metadata_json", None) or {}
        delivery = metadata.get("delivery") if isinstance(metadata, dict) else None
        readiness = delivery.get("delivery_readiness") if isinstance(delivery, dict) else None
        return isinstance(readiness, dict) and readiness.get("ready") is False

    async def _count_unresolved_document_comments(
        self,
        document_id: UUID,
        tenant_id: UUID,
    ) -> int:
        from app.domains.collaboration.models import DocumentComment

        count = await self.db.scalar(
            select(func.count(DocumentComment.id)).where(
                DocumentComment.document_id == document_id,
                DocumentComment.tenant_id == tenant_id,
                DocumentComment.resolved.is_(False),
                DocumentComment.deleted_at.is_(None),
            )
        )
        return int(count or 0)

    async def _count_failed_document_quality_checks(
        self,
        document_id: UUID,
        tenant_id: UUID,
    ) -> int:
        from app.domains.documents.models import QualityResult

        result = await self.db.execute(
            select(QualityResult).where(
                QualityResult.document_id == document_id,
                QualityResult.tenant_id == tenant_id,
            )
        )
        failed_count = 0
        for quality_result in result.scalars().all():
            issues = quality_result.issues_json or {}
            status = str(issues.get("status") or "").lower() if isinstance(issues, dict) else ""
            if quality_result.score < 70 or status in {"blocked", "error", "fail", "failed"}:
                failed_count += 1
        return failed_count

    async def create_export_job(
        self,
        tenant_id: UUID,
        project_id: UUID,
        document_id: UUID | None,
        template_id: UUID | None,
        export_type: str,
        created_by: UUID | None = None,
    ) -> ExportJob:
        """Create a new export job.

        Args:
            tenant_id: Tenant UUID
            project_id: Project UUID
            document_id: Document UUID to export
            template_id: Optional template UUID for formatting
            export_type: Type of export (word, markdown, pptx)
            created_by: User ID of requester

        Returns:
            Created ExportJob
        """
        if created_by is None:
            raise ValueError("created_by is required")

        job = ExportJob(
            tenant_id=tenant_id,
            project_id=project_id,
            document_id=document_id,
            template_id=template_id,
            export_type=export_type,
            status=ExportStatus.PENDING.value,
            created_by=created_by,
        )
        self.db.add(job)
        await self.db.flush()
        await self.db.refresh(job)
        return job

    async def get_job(
        self,
        job_id: UUID,
        tenant_id: UUID | None = None,
    ) -> ExportJob | None:
        """Get export job by ID.

        Args:
            job_id: ExportJob UUID
            tenant_id: Optional tenant filter

        Returns:
            ExportJob if found, None otherwise
        """
        query = select(ExportJob).where(ExportJob.id == job_id)
        if tenant_id is not None:
            query = query.where(ExportJob.tenant_id == tenant_id)

        result = await self.db.execute(query)
        return result.scalar_one_or_none()

    async def get_job_status(
        self,
        job_id: UUID,
        tenant_id: UUID | None = None,
    ) -> ExportJob | None:
        """Get export job status.

        Args:
            job_id: ExportJob UUID
            tenant_id: Optional tenant filter

        Returns:
            ExportJob if found, None otherwise
        """
        return await self.get_job(job_id, tenant_id)

    async def list_jobs(self, tenant_id: UUID) -> list[ExportJob]:
        """List all export jobs for a tenant.

        Args:
            tenant_id: Tenant UUID

        Returns:
            List of ExportJobs
        """
        query = (
            select(ExportJob)
            .where(ExportJob.tenant_id == tenant_id)
            .order_by(ExportJob.created_at.desc())
        )
        result = await self.db.execute(query)
        return list(result.scalars().all())

    async def export_project_package(
        self,
        project_id: UUID,
        tenant_id: UUID,
        document_ids: list[UUID] | None = None,
        title: str | None = None,
        include_drafts: bool = False,
        include_manifest: bool = True,
        formats: list[str] | None = None,
        include_audit: bool = False,
        watermark: str | None = None,
        variables: dict[str, Any] | None = None,
        created_by: UUID | None = None,
        user_role_id: UUID | None = None,
    ) -> ExportJob:
        """Export a project delivery package as selected downloadable artifacts."""
        from app.domains.documents.models import Document, DocumentStatus
        from app.models.projects import Project

        project_result = await self.db.execute(
            select(Project).where(
                Project.id == project_id,
                Project.tenant_id == tenant_id,
                Project.deleted_at.is_(None),
            )
        )
        project = project_result.scalar_one_or_none()
        if not project:
            raise ValueError(f"Project not found: {project_id}")

        query = select(Document).where(
            Document.project_id == project_id,
            Document.tenant_id == tenant_id,
            Document.deleted_at.is_(None),
        )

        requested_ids = set(document_ids or [])
        if requested_ids:
            query = query.where(Document.id.in_(requested_ids))
        elif not include_drafts:
            query = query.where(Document.status == DocumentStatus.PUBLISHED.value)

        result = await self.db.execute(query)
        documents = list(result.scalars().all())

        if requested_ids:
            found_ids = {doc.id for doc in documents}
            missing_ids = requested_ids - found_ids
            if missing_ids:
                missing = ", ".join(str(item) for item in sorted(missing_ids, key=str))
                raise ValueError(f"Documents not found in project: {missing}")

        placeholder_documents = [
            doc for doc in documents
            if str((doc.metadata_json or {}).get("generation_status") or "").lower()
            in {"placeholder", "partial", "failed"}
        ]
        if placeholder_documents:
            names = ", ".join(doc.title for doc in placeholder_documents)
            raise ValueError(
                "Cannot export placeholder documents or other non-generated AI documents: "
                f"{names}"
            )

        readiness_blocked_documents = [
            doc for doc in documents
            if self._document_delivery_readiness_failed(doc)
        ]
        if readiness_blocked_documents:
            names = ", ".join(doc.title for doc in readiness_blocked_documents)
            raise ValueError(f"Cannot export documents with export readiness blockers: {names}")

        unresolved_comment_documents: list[Any] = []
        failed_quality_documents: list[Any] = []
        for document in documents:
            if await self._count_unresolved_document_comments(document.id, tenant_id):
                unresolved_comment_documents.append(document)
            if await self._count_failed_document_quality_checks(document.id, tenant_id):
                failed_quality_documents.append(document)
        if unresolved_comment_documents:
            names = ", ".join(doc.title for doc in unresolved_comment_documents)
            raise ValueError(f"Cannot export documents with unresolved comments: {names}")
        if failed_quality_documents:
            names = ", ".join(doc.title for doc in failed_quality_documents)
            raise ValueError(f"Cannot export documents with failed quality checks: {names}")

        if not documents:
            raise ValueError("No exportable project documents found")

        if user_role_id:
            documents = [
                await self._filter_document_by_permissions(doc, user_role_id, tenant_id)
                for doc in documents
            ]

        documents = sorted(
            documents,
            key=lambda doc: (
                self._document_package_order(getattr(doc, "doc_type", "")),
                (getattr(doc, "title", "") or "").lower(),
            ),
        )

        job = await self.create_export_job(
            tenant_id=tenant_id,
            project_id=project_id,
            document_id=None,
            template_id=None,
            export_type=ExportType.PROJECT_PACKAGE.value,
            created_by=created_by,
        )
        await self.update_job_status(job.id, ExportStatus.PROCESSING.value)

        try:
            package_title = title or f"{project.name} 交付包"
            package_formats = self._normalize_package_formats(formats)
            markdown_content = self._build_project_package_markdown(
                project_name=project.name,
                title=package_title,
                documents=documents,
                include_manifest=include_manifest,
                include_audit=include_audit,
                watermark=watermark,
                variables=variables or {},
            )
            safe_title = self._safe_filename(package_title)
            first_handle: StorageHandle | None = None

            for package_format in package_formats:
                filename, content, content_type = self._build_project_package_artifact(
                    package_format=package_format,
                    safe_title=safe_title,
                    package_title=package_title,
                    project_name=project.name,
                    documents=documents,
                    markdown_content=markdown_content,
                    include_manifest=include_manifest,
                    include_audit=include_audit,
                    watermark=watermark,
                    variables=variables or {},
                )
                handle = await self._upload_project_package_artifact(
                    tenant_id=tenant_id,
                    project_id=project_id,
                    job_id=job.id,
                    filename=filename,
                    content=content,
                    content_type=content_type,
                )
                if first_handle is None:
                    first_handle = handle

            if first_handle is None:
                raise ValueError("No package artifacts generated")

            await self.update_job_status(
                job.id,
                ExportStatus.COMPLETED.value,
                output_path=first_handle.path,
                file_hash=first_handle.hash,
            )
        except Exception as exc:
            await self.update_job_status(
                job.id,
                ExportStatus.FAILED.value,
                error_message=str(exc),
            )
            raise

        return job

    def _build_project_package_markdown(
        self,
        project_name: str,
        title: str,
        documents: list[Any],
        include_manifest: bool,
        include_audit: bool,
        watermark: str | None,
        variables: dict[str, Any],
    ) -> str:
        """Build a single Markdown delivery package from project documents."""
        lines = [
            f"# {title}",
            "",
            f"- 项目：{project_name}",
            f"- 文档数量：{len(documents)}",
            f"- 生成时间：{datetime.now(timezone.utc).isoformat()}",
            "",
        ]

        if watermark:
            lines.extend(["## Watermark", "", watermark, ""])

        if include_manifest:
            lines.extend(["## 交付清单", ""])
            for index, document in enumerate(documents, start=1):
                lines.append(
                    f"{index}. {document.title} "
                    f"({document.doc_type}, v{document.version}, {document.status})"
                )
            lines.append("")

        if include_audit:
            lines.extend([
                "## Audit checklist",
                "",
                "- Document count verified.",
                "- Placeholder documents blocked before export.",
                "- Selected documents sorted by delivery order.",
                "- Artifact generation completed for requested formats.",
                "",
            ])

        for index, document in enumerate(documents, start=1):
            lines.extend([
                "---",
                "",
                f"## {index}. {document.title}",
                "",
                f"- 类型：{document.doc_type}",
                f"- 状态：{document.status}",
                f"- 版本：{document.version}",
                "",
                self._substitute_variables(document.content or "", variables),
                "",
            ])

        return "\n".join(lines).strip() + "\n"

    def _normalize_package_formats(self, formats: list[str] | None) -> list[str]:
        requested = formats or ["markdown"]
        supported = {"markdown", "word", "pptx"}
        normalized: list[str] = []
        for item in requested:
            value = str(item).strip().lower()
            if not value:
                continue
            if value not in supported:
                raise ValueError(f"Unsupported package export format: {item}")
            if value not in normalized:
                normalized.append(value)
        if not normalized:
            raise ValueError("At least one package export format is required")
        return normalized

    def _build_project_package_artifact(
        self,
        *,
        package_format: str,
        safe_title: str,
        package_title: str,
        project_name: str,
        documents: list[Any],
        markdown_content: str,
        include_manifest: bool,
        include_audit: bool,
        watermark: str | None,
        variables: dict[str, Any],
    ) -> tuple[str, bytes, str]:
        if package_format == "markdown":
            return (
                f"{safe_title}.project-package.md",
                markdown_content.encode("utf-8"),
                "text/markdown",
            )
        if package_format == "word":
            return (
                f"{safe_title}.project-package.docx",
                self._build_project_package_docx(
                    project_name=project_name,
                    title=package_title,
                    documents=documents,
                    include_manifest=include_manifest,
                    include_audit=include_audit,
                    watermark=watermark,
                    variables=variables,
                ),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        return (
            f"{safe_title}.project-package.pptx",
            self._build_project_package_pptx(
                project_name=project_name,
                title=package_title,
                documents=documents,
                include_audit=include_audit,
                watermark=watermark,
                variables=variables,
            ),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    async def _upload_project_package_artifact(
        self,
        *,
        tenant_id: UUID,
        project_id: UUID,
        job_id: UUID,
        filename: str,
        content: bytes,
        content_type: str,
    ) -> StorageHandle:
        storage = get_storage_provider()
        handle = await storage.upload(
            tenant_id=str(tenant_id),
            project_id=str(project_id),
            filename=filename,
            content=content,
            content_type=content_type,
        )
        await self.create_artifact(
            tenant_id=tenant_id,
            job_id=job_id,
            filename=filename,
            content_type=content_type,
            file_size=len(content),
            storage_path=handle.path,
            file_hash=handle.hash,
        )
        return handle

    def _build_project_package_docx(
        self,
        *,
        project_name: str,
        title: str,
        documents: list[Any],
        include_manifest: bool,
        include_audit: bool,
        watermark: str | None,
        variables: dict[str, Any],
    ) -> bytes:
        from docx import Document as WordDocument

        docx = WordDocument()
        docx.add_heading(title, level=0)
        docx.add_paragraph(f"Project: {project_name}")
        docx.add_paragraph(f"Generated at: {datetime.now(timezone.utc).isoformat()}")
        if watermark:
            docx.add_paragraph(f"Watermark: {watermark}")
        if include_manifest:
            docx.add_heading("Delivery manifest", level=1)
            for index, document in enumerate(documents, start=1):
                docx.add_paragraph(
                    f"{index}. {document.title} ({document.doc_type}, v{document.version}, {document.status})"
                )
        if include_audit:
            docx.add_heading("Audit checklist", level=1)
            for item in [
                "Document count verified.",
                "Placeholder documents blocked before export.",
                "Selected documents sorted by delivery order.",
                "Artifact generation completed for requested formats.",
            ]:
                docx.add_paragraph(item, style="List Bullet")
        for index, document in enumerate(documents, start=1):
            docx.add_heading(f"{index}. {document.title}", level=1)
            docx.add_paragraph(f"Type: {document.doc_type}")
            docx.add_paragraph(f"Status: {document.status}")
            docx.add_paragraph(f"Version: {document.version}")
            docx.add_paragraph(self._substitute_variables(document.content or "", variables))
        buffer = io.BytesIO()
        docx.save(buffer)
        return buffer.getvalue()

    def _build_project_package_pptx(
        self,
        *,
        project_name: str,
        title: str,
        documents: list[Any],
        include_audit: bool,
        watermark: str | None,
        variables: dict[str, Any],
    ) -> bytes:
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        title_slide.placeholders[1].text = f"{project_name}\n{watermark or 'Delivery package'}"

        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        summary_slide.shapes.title.text = "Delivery manifest"
        body = summary_slide.shapes.placeholders[1].text_frame
        body.text = f"Documents: {len(documents)}"
        if include_audit:
            paragraph = body.add_paragraph()
            paragraph.text = "Audit checklist included"

        for index, document in enumerate(documents, start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"{index}. {document.title}"
            text_frame = slide.shapes.placeholders[1].text_frame
            text_frame.text = f"{document.doc_type} | {document.status} | v{document.version}"
            paragraph = text_frame.add_paragraph()
            paragraph.text = self._substitute_variables(document.content or "", variables)[:700]

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _document_package_order(self, doc_type: str) -> int:
        order = {
            "urs": 10,
            "brd": 20,
            "prd": 30,
            "user_story": 40,
            "detailed_design": 50,
            "interface": 60,
            "data_dictionary": 70,
            "test_case": 80,
        }
        return order.get(doc_type, 999)

    def _safe_filename(self, value: str) -> str:
        filename = re.sub(r"[^\w\u4e00-\u9fff -]+", "_", value).strip()
        return re.sub(r"\s+", "-", filename) or "project-package"

    def _substitute_variables(self, content: str, variables: dict[str, Any]) -> str:
        return substitute_placeholders(content, variables)

    async def update_job_status(
        self,
        job_id: UUID,
        status: str,
        output_path: str | None = None,
        file_hash: str | None = None,
        error_message: str | None = None,
    ) -> ExportJob | None:
        """Update export job status.

        Args:
            job_id: ExportJob UUID
            status: New status
            output_path: Optional output file path
            file_hash: Optional file hash
            error_message: Optional error message

        Returns:
            Updated ExportJob if found, None otherwise
        """
        job = await self.get_job(job_id)
        if not job:
            return None

        job.status = status
        if output_path is not None:
            job.output_path = output_path
        if file_hash is not None:
            job.file_hash = file_hash
        if error_message is not None:
            job.error_message = error_message
        if status in (ExportStatus.COMPLETED.value, ExportStatus.FAILED.value):
            job.completed_at = datetime.now(timezone.utc)

        await self.db.flush()
        await self.db.refresh(job)
        return job

    async def get_artifact(
        self,
        artifact_id: UUID,
        tenant_id: UUID | None = None,
    ) -> ExportArtifact | None:
        """Get export artifact by ID.

        Args:
            artifact_id: ExportArtifact UUID
            tenant_id: Optional tenant filter

        Returns:
            ExportArtifact if found, None otherwise
        """
        query = select(ExportArtifact).where(ExportArtifact.id == artifact_id)
        if tenant_id is not None:
            query = query.where(ExportArtifact.tenant_id == tenant_id)

        result = await self.db.execute(query)
        return result.scalar_one_or_none()

    async def create_artifact(
        self,
        tenant_id: UUID,
        job_id: UUID,
        filename: str,
        content_type: str,
        file_size: int,
        storage_path: str,
        file_hash: str | None = None,
    ) -> ExportArtifact:
        """Create an export artifact record.

        Args:
            tenant_id: Tenant UUID
            job_id: Parent job UUID
            filename: Original filename
            content_type: MIME type
            file_size: File size in bytes
            storage_path: Path in storage
            file_hash: Optional SHA256 hash

        Returns:
            Created ExportArtifact
        """
        artifact = ExportArtifact(
            tenant_id=tenant_id,
            job_id=job_id,
            filename=filename,
            content_type=content_type,
            file_size=file_size,
            storage_path=storage_path,
            file_hash=file_hash,
        )
        self.db.add(artifact)
        await self.db.flush()
        await self.db.refresh(artifact)
        return artifact

    async def _filter_document_by_permissions(
        self,
        document: Any,
        role_id: UUID,
        tenant_id: UUID,
    ) -> Any:
        """Filter document fields based on field permissions.

        Args:
            document: Document to filter
            role_id: Role UUID for permission lookup
            tenant_id: Tenant UUID

        Returns:
            Document with filtered fields (or dict with filtered content)
        """
        from app.domains.identity.service import FieldPermissionService

        field_perm_service = FieldPermissionService(self.db)
        permissions = await field_perm_service.get_field_permissions_for_role(
            role_id=role_id,
            resource_type="document",
            tenant_id=tenant_id,
        )

        if not permissions:
            return document

        # Build permission map
        perm_map: dict[str, str] = {}
        for perm in permissions:
            perm_map[perm.field_name] = perm.permission

        # Get document content and filter based on permissions
        content = document.content or ""
        metadata = document.metadata_json or {}

        # If content field has no read permission, redact it
        if perm_map.get("content") == "none":
            content = "[Content redacted due to field-level permissions]"

        # Filter metadata fields
        filtered_metadata = {}
        for field, value in metadata.items():
            perm = perm_map.get(field)
            if perm is None or perm in ("read", "write"):
                filtered_metadata[field] = value
            # "none" permission means field is excluded

        # Create a new document-like object with filtered content
        # We use a simple namespace object to hold filtered values
        class FilteredDocument:
            def __init__(self, doc, filtered_content, filtered_metadata):
                self.id = doc.id
                self.project_id = doc.project_id
                self.doc_type = doc.doc_type
                self.title = doc.title
                self.content = filtered_content
                self.status = doc.status
                self.version = doc.version
                self.parent_document_id = doc.parent_document_id
                self.created_by = doc.created_by
                self.approved_by = doc.approved_by
                self.quality_score = doc.quality_score
                self.metadata_json = filtered_metadata
                self.created_at = doc.created_at
                self.updated_at = doc.updated_at

        return FilteredDocument(document, content, filtered_metadata)

    async def export_word(
        self,
        document_id: UUID,
        template_id: UUID | None,
        tenant_id: UUID,
        variables: dict[str, Any] | None = None,
        title: str | None = None,
        created_by: UUID | None = None,
        user_role_id: UUID | None = None,
    ) -> ExportJob:
        """Export document to Word format.

        Args:
            document_id: Document UUID to export
            template_id: Optional template UUID
            tenant_id: Tenant UUID
            variables: Additional variables for substitution
            title: Optional title override
            created_by: User ID of requester
            user_role_id: Optional role ID for field permission filtering

        Returns:
            Created ExportJob
        """
        # Get document
        from app.domains.documents.service import DocumentService
        doc_service = DocumentService(self.db)
        document = await doc_service.get_document(document_id, tenant_id)
        if not document:
            raise ValueError(f"Document not found: {document_id}")

        # Block export of placeholder documents
        metadata = document.metadata_json or {}
        if metadata.get("generation_status") == "placeholder":
            raise ValueError(
                f"Cannot export placeholder document. Document must be regenerated with LLM before export."
            )

        # Apply field permission filtering if role_id provided
        if user_role_id:
            document = await self._filter_document_by_permissions(
                document, user_role_id, tenant_id
            )

        # Get project_id from document
        project_id = document.project_id

        # Create job
        job = await self.create_export_job(
            tenant_id=tenant_id,
            project_id=project_id,
            document_id=document_id,
            template_id=template_id,
            export_type=ExportType.WORD.value,
            created_by=created_by,
        )

        # Update status to processing
        await self.update_job_status(job.id, ExportStatus.PROCESSING.value)

        # Perform export
        try:
            exporter = WordExporter(self.db, tenant_id)
            result = await exporter.export(
                document=document,
                template_id=template_id,
                variables=variables or {},
                title=title,
            )

            # Store artifact
            storage = get_storage_provider()
            handle = await storage.upload(
                tenant_id=str(tenant_id),
                project_id=str(project_id),
                filename=result["filename"],
                content=result["content"],
                content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )

            # Create artifact record
            await self.create_artifact(
                tenant_id=tenant_id,
                job_id=job.id,
                filename=result["filename"],
                content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                file_size=len(result["content"]),
                storage_path=handle.path,
                file_hash=handle.hash,
            )

            # Update job status
            await self.update_job_status(
                job.id,
                ExportStatus.COMPLETED.value,
                output_path=handle.path,
                file_hash=handle.hash,
            )
        except Exception as e:
            await self.update_job_status(
                job.id,
                ExportStatus.FAILED.value,
                error_message=str(e),
            )
            raise

        return job

    async def export_markdown(
        self,
        document_id: UUID,
        tenant_id: UUID,
        variables: dict[str, Any] | None = None,
        title: str | None = None,
        created_by: UUID | None = None,
        user_role_id: UUID | None = None,
    ) -> ExportJob:
        """Export document to Markdown format.

        Args:
            document_id: Document UUID to export
            tenant_id: Tenant UUID
            variables: Additional variables for substitution
            title: Optional title override
            created_by: User ID of requester
            user_role_id: Optional role ID for field permission filtering

        Returns:
            Created ExportJob
        """
        from app.domains.documents.service import DocumentService
        doc_service = DocumentService(self.db)
        document = await doc_service.get_document(document_id, tenant_id)
        if not document:
            raise ValueError(f"Document not found: {document_id}")

        # Block export of placeholder documents
        metadata = document.metadata_json or {}
        if metadata.get("generation_status") == "placeholder":
            raise ValueError(
                f"Cannot export placeholder document. Document must be regenerated with LLM before export."
            )

        # Apply field permission filtering if role_id provided
        if user_role_id:
            document = await self._filter_document_by_permissions(
                document, user_role_id, tenant_id
            )

        project_id = document.project_id

        job = await self.create_export_job(
            tenant_id=tenant_id,
            project_id=project_id,
            document_id=document_id,
            template_id=None,
            export_type=ExportType.MARKDOWN.value,
            created_by=created_by,
        )

        await self.update_job_status(job.id, ExportStatus.PROCESSING.value)

        try:
            exporter = MarkdownExporter(self.db, tenant_id)
            result = await exporter.export(
                document=document,
                variables=variables or {},
                title=title,
            )

            storage = get_storage_provider()
            handle = await storage.upload(
                tenant_id=str(tenant_id),
                project_id=str(project_id),
                filename=result["filename"],
                content=result["content"],
                content_type="text/markdown",
            )

            await self.create_artifact(
                tenant_id=tenant_id,
                job_id=job.id,
                filename=result["filename"],
                content_type="text/markdown",
                file_size=len(result["content"]),
                storage_path=handle.path,
                file_hash=handle.hash,
            )

            await self.update_job_status(
                job.id,
                ExportStatus.COMPLETED.value,
                output_path=handle.path,
                file_hash=handle.hash,
            )
        except Exception as e:
            await self.update_job_status(
                job.id,
                ExportStatus.FAILED.value,
                error_message=str(e),
            )
            raise

        return job

    async def export_pptx(
        self,
        document_id: UUID,
        template_id: UUID | None,
        tenant_id: UUID,
        variables: dict[str, Any] | None = None,
        title: str | None = None,
        created_by: UUID | None = None,
        user_role_id: UUID | None = None,
    ) -> ExportJob:
        """Export document to PPTX format.

        Args:
            document_id: Document UUID to export
            template_id: Optional template UUID
            tenant_id: Tenant UUID
            variables: Additional variables for substitution
            title: Optional title override
            created_by: User ID of requester
            user_role_id: Optional role ID for field permission filtering

        Returns:
            Created ExportJob
        """
        from app.domains.documents.service import DocumentService
        doc_service = DocumentService(self.db)
        document = await doc_service.get_document(document_id, tenant_id)
        if not document:
            raise ValueError(f"Document not found: {document_id}")

        # Block export of placeholder documents
        metadata = document.metadata_json or {}
        if metadata.get("generation_status") == "placeholder":
            raise ValueError(
                f"Cannot export placeholder document. Document must be regenerated with LLM before export."
            )

        # Apply field permission filtering if role_id provided
        if user_role_id:
            document = await self._filter_document_by_permissions(
                document, user_role_id, tenant_id
            )

        project_id = document.project_id

        job = await self.create_export_job(
            tenant_id=tenant_id,
            project_id=project_id,
            document_id=document_id,
            template_id=template_id,
            export_type=ExportType.PPTX.value,
            created_by=created_by,
        )

        await self.update_job_status(job.id, ExportStatus.PROCESSING.value)

        try:
            exporter = PPTXExporter(self.db, tenant_id)
            result = await exporter.export(
                document=document,
                template_id=template_id,
                variables=variables or {},
                title=title,
            )

            storage = get_storage_provider()
            handle = await storage.upload(
                tenant_id=str(tenant_id),
                project_id=str(project_id),
                filename=result["filename"],
                content=result["content"],
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

            await self.create_artifact(
                tenant_id=tenant_id,
                job_id=job.id,
                filename=result["filename"],
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                file_size=len(result["content"]),
                storage_path=handle.path,
                file_hash=handle.hash,
            )

            await self.update_job_status(
                job.id,
                ExportStatus.COMPLETED.value,
                output_path=handle.path,
                file_hash=handle.hash,
            )
        except Exception as e:
            await self.update_job_status(
                job.id,
                ExportStatus.FAILED.value,
                error_message=str(e),
            )
            raise

        return job


class WordExporter:
    """Exporter for Word (.docx) format using python-docx.

    Supports template variable substitution with {{variable}} syntax.
    """

    def __init__(self, db: AsyncSession, tenant_id: UUID):
        """Initialize Word exporter.

        Args:
            db: Async database session
            tenant_id: Tenant UUID
        """
        self.db = db
        self.tenant_id = tenant_id

    async def export(
        self,
        document: Any,
        template_id: UUID | None,
        variables: dict[str, Any],
        title: str | None = None,
    ) -> dict[str, Any]:
        """Export document to Word format.

        Args:
            document: Document to export
            template_id: Optional template UUID
            variables: Variables for substitution
            title: Optional title override

        Returns:
            Dictionary with 'filename' and 'content' (bytes)
        """
        from docx import Document
        from docx.shared import Pt
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

        # Get template content if provided
        template_content = None
        if template_id:
            from app.domains.templates.service import TemplateService
            svc = TemplateService(self.db)
            active_version = await svc.get_active_version(template_id, self.tenant_id)
            if active_version and active_version.content:
                template_content = active_version.content

        # Create document
        doc = Document()

        # Set title
        doc_title = title or document.title
        heading = doc.add_heading(doc_title, 0)
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        # Get content and apply variable substitution
        content = document.content or ""
        content = self._substitute_variables(content, variables)

        # Add paragraphs
        for line in content.split("\n"):
            if line.startswith("# "):
                doc.add_heading(line[2:], level=1)
            elif line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.startswith("### "):
                doc.add_heading(line[4:], level=3)
            elif line.startswith("- ") or line.startswith("* "):
                doc.add_paragraph(line[2:], style="List Bullet")
            elif line.strip():
                doc.add_paragraph(line)

        # Save to bytes
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        content_bytes = buffer.read()

        # Generate filename
        safe_title = "".join(c if c.isalnum() or c in " -_" else "_" for c in doc_title)
        filename = f"{safe_title}.docx"

        return {
            "filename": filename,
            "content": content_bytes,
        }

    def _substitute_variables(self, content: str, variables: dict[str, Any]) -> str:
        """Substitute {{variable}} placeholders with values.

        Args:
            content: Content with placeholders
            variables: Variable values

        Returns:
            Content with substituted values
        """
        return substitute_placeholders(content, variables)


class MarkdownExporter:
    """Exporter for Markdown format.

    Simple text generation with markdown formatting.
    """

    def __init__(self, db: AsyncSession, tenant_id: UUID):
        """Initialize Markdown exporter.

        Args:
            db: Async database session
            tenant_id: Tenant UUID
        """
        self.db = db
        self.tenant_id = tenant_id

    async def export(
        self,
        document: Any,
        variables: dict[str, Any],
        title: str | None = None,
    ) -> dict[str, Any]:
        """Export document to Markdown format.

        Args:
            document: Document to export
            variables: Variables for substitution
            title: Optional title override

        Returns:
            Dictionary with 'filename' and 'content' (bytes)
        """
        # Get document title
        doc_title = title or document.title

        # Get content and apply variable substitution
        content = document.content or ""
        content = self._substitute_variables(content, variables)

        # Format as markdown
        markdown_lines = [
            f"# {doc_title}",
            "",
            content,
            "",
        ]

        markdown_content = "\n".join(markdown_lines)
        content_bytes = markdown_content.encode("utf-8")

        # Generate filename
        safe_title = "".join(c if c.isalnum() or c in " -_" else "_" for c in doc_title)
        filename = f"{safe_title}.md"

        return {
            "filename": filename,
            "content": content_bytes,
        }

    def _substitute_variables(self, content: str, variables: dict[str, Any]) -> str:
        """Substitute {{variable}} placeholders with values.

        Args:
            content: Content with placeholders
            variables: Variable values

        Returns:
            Content with substituted values
        """
        return substitute_placeholders(content, variables)


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import zipfile
from dataclasses import dataclass, field
from typing import Any


@dataclass
class PlaceholderInfo:
    """Information about a placeholder in a slide."""
    variable_name: str
    shape_id: int
    shape_name: str
    placeholder_type: str | None  # title, body, chart, table, etc.
    text: str
    position: tuple[int, int]  # (paragraph_index, run_index)


@dataclass
class SlideStructure:
    """Structure information for a slide."""
    slide_index: int
    layout_name: str
    layout_type: str  # title, content, two_column, blank, etc.
    placeholders: list[PlaceholderInfo] = field(default_factory=list)
    shapes_with_placeholders: list[str] = field(default_factory=list)


@dataclass
class ContentMapping:
    """Mapping of content to slide placeholders."""
    slide_index: int
    placeholder_type: str
    variable_name: str
    content: str
    style_info: dict[str, Any] = field(default_factory=dict)


class PPTXExporter:
    """Exporter for PowerPoint (.pptx) format using python-pptx.

    STRICT placeholder parsing: only supported {{variable_name}} or
    {{中文变量}} tokens in plain text boxes are recognized. Page type
    detection: title, content, chart, table.
    """

    # Pattern for strict placeholder extraction
    PLACEHOLDER_PATTERN = SHARED_PLACEHOLDER_PATTERN.pattern

    def __init__(self, db: AsyncSession, tenant_id: UUID):
        """Initialize PPTX exporter.

        Args:
            db: Async database session
            tenant_id: Tenant UUID
        """
        self.db = db
        self.tenant_id = tenant_id

    async def export(
        self,
        document: Any,
        template_id: UUID | None,
        variables: dict[str, Any],
        title: str | None = None,
    ) -> dict[str, Any]:
        """Export document to PPTX format.

        Args:
            document: Document to export
            template_id: Optional template UUID
            variables: Variables for substitution
            title: Optional title override

        Returns:
            Dictionary with 'filename' and 'content' (bytes)
        """
        # Get template content if provided
        template_content = None
        if template_id:
            from app.domains.templates.service import TemplateService
            svc = TemplateService(self.db)
            active_version = await svc.get_active_version(template_id, self.tenant_id)
            if active_version and active_version.content:
                template_content = active_version.content

        # Create presentation
        prs = Presentation()

        # Get document content
        doc_title = title or document.title
        content = document.content or ""

        # If we have a template, use it as base
        if template_content:
            # Load from template bytes
            buffer = io.BytesIO(template_content)
            prs = Presentation(buffer)

            # Find and process placeholders in existing slides
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if self._contains_placeholder(run.text):
                                    # Extract and substitute
                                    text = run.text
                                    text = self._substitute_variables(text, variables)
                                    run.text = text
                    elif hasattr(shape, "text") and isinstance(shape.text, str):
                        if self._contains_placeholder(shape.text):
                            shape.text = self._substitute_variables(shape.text, variables)
        else:
            # Generate slides from document content
            self._create_slides_from_content(prs, doc_title, content, variables)

        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        content_bytes = buffer.read()

        # Generate filename
        safe_title = "".join(c if c.isalnum() or c in " -_" else "_" for c in doc_title)
        filename = f"{safe_title}.pptx"

        return {
            "filename": filename,
            "content": content_bytes,
        }

    def _contains_placeholder(self, text: str) -> bool:
        """Check if text contains a valid placeholder.

        Args:
            text: Text to check

        Returns:
            True if text contains {{variable}} pattern
        """
        return contains_placeholder(text)

    def _extract_placeholders(self, text: str) -> list[str]:
        """Extract variable names from placeholder text.

        Args:
            text: Text containing {{variable_name}} patterns

        Returns:
            List of variable names
        """
        return extract_placeholders(text)

    def _substitute_variables(self, content: str, variables: dict[str, Any]) -> str:
        """Substitute {{variable}} placeholders with provided values.

        Args:
            content: Content with placeholders
            variables: Variable values

        Returns:
            Content with substituted values
        """
        return substitute_placeholders(content, variables)

    def _create_slides_from_content(
        self,
        prs: Any,
        title: str,
        content: str,
        variables: dict[str, Any],
    ) -> None:
        """Create slides from document content.

        Args:
            prs: Presentation object
            title: Document title
            content: Document content
            variables: Variables for substitution
        """
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = self._substitute_variables(title, variables)

        # Content slides
        content_lines = content.split("\n")
        current_content = []

        for line in content_lines:
            line = line.strip()
            if not line:
                continue

            if line.startswith("# ") or line.startswith("## "):
                # Save previous content slide
                if current_content:
                    self._add_content_slide(prs, current_content, variables)
                    current_content = []

                # Add heading as new slide
                heading_text = line.lstrip("# ").strip()
                self._add_title_slide(prs, heading_text, variables)
            elif line.startswith("- ") or line.startswith("* "):
                current_content.append(("bullet", line[2:].strip()))
            else:
                current_content.append(("text", line))

        # Add final content slide
        if current_content:
            self._add_content_slide(prs, current_content, variables)

    def _add_title_slide(
        self,
        prs: Any,
        title: str,
        variables: dict[str, Any],
    ) -> None:
        """Add a title slide.

        Args:
            prs: Presentation object
            title: Slide title
            variables: Variables for substitution
        """
        try:
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = self._substitute_variables(title, variables)
        except Exception:
            # Fallback to blank layout
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(blank_layout)

    def _add_content_slide(
        self,
        prs: Any,
        content_items: list[tuple[str, str]],
        variables: dict[str, Any],
    ) -> None:
        """Add a content slide with bullet points.

        Args:
            prs: Presentation object
            content_items: List of (type, text) tuples
            variables: Variables for substitution
        """
        try:
            content_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_layout)
        except Exception:
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(blank_layout)

        # Try to find a content placeholder
        body_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "placeholder_format") and hasattr(shape, "text_frame"):
                body_shape = shape
                break

        if body_shape and hasattr(body_shape, "text_frame"):
            tf = body_shape.text_frame
            tf.clear()

            for item_type, text in content_items:
                p = tf.add_paragraph()
                p.text = self._substitute_variables(text, variables)
                if item_type == "bullet":
                    p.level = 0
                else:
                    p.level = 0
        else:
            # Fallback: add text box
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.clear()

            for item_type, text in content_items:
                p = tf.add_paragraph()
                p.text = self._substitute_variables(text, variables)
                if item_type == "bullet":
                    p.level = 0

    # =========================================================================
    # L2: Structure Detection and Content Fill Methods
    # =========================================================================

    def _parse_pptx_structure(self, pptx_path: str) -> dict[str, Any]:
        """Parse PPTX ZIP to find all slides, shapes, and placeholders.

        Args:
            pptx_path: Path to the PPTX file

        Returns:
            Dictionary with slide structures and metadata
        """
        structure: dict[str, Any] = {
            "slides": [],
            "layouts": [],
            "slide_count": 0,
        }

        try:
            with zipfile.ZipFile(pptx_path, "r") as zf:
                # Parse [Content_Types].xml to find all slides
                content_types_xml = zf.read("[Content_Types].xml")
                ct_root = etree.fromstring(content_types_xml)

                # Find all slide references
                slide_refs = []
                for elem in ct_root.iter():
                    if elem.get("PartName", "").startswith("/ppt/slides/slide"):
                        slide_refs.append(elem.get("PartName"))

                structure["slide_count"] = len(slide_refs)

                # Parse each slide
                for slide_ref in sorted(slide_refs):
                    slide_xml = zf.read(slide_ref.lstrip("/"))
                    slide_root = etree.fromstring(slide_xml)

                    # Extract slide name
                    slide_name = slide_ref.split("/")[-1]

                    # Get slide relationships
                    slide_rel_path = slide_ref.replace(".xml", ".xml.rels")
                    try:
                        slide_rel_xml = zf.read(slide_rel_path.lstrip("/"))
                        rel_root = etree.fromstring(slide_rel_xml)
                    except KeyError:
                        rel_root = None

                    structure["slides"].append({
                        "name": slide_name,
                        "path": slide_ref,
                        "xml": etree.tostring(slide_root, pretty_print=True).decode("utf-8"),
                    })

                # Parse layouts
                layouts_path = "ppt/slideLayouts/_rels"
                try:
                    layouts_xml = zf.read("ppt/slideLayouts/slideLayouts.xml")
                    layouts_root = etree.fromstring(layouts_xml)
                    for elem in layouts_root.iter():
                        if elem.get("Type", "").endswith("slideLayout"):
                            structure["layouts"].append(elem.get("id", "unknown"))
                except KeyError:
                    pass

        except Exception as e:
            structure["error"] = str(e)

        return structure

    def _detect_slide_layouts(self, prs: Presentation) -> list[dict[str, Any]]:
        """Identify layout types for all slides in a presentation.

        Args:
            prs: Presentation object

        Returns:
            List of slide layout information
        """
        layouts: list[dict[str, Any]] = []

        for idx, slide in enumerate(prs.slides):
            layout_info: dict[str, Any] = {
                "slide_index": idx,
                "layout_name": "unknown",
                "layout_type": "unknown",
                "placeholder_count": 0,
                "has_title": False,
                "has_body": False,
                "has_chart": False,
                "has_table": False,
            }

            # Get layout name
            try:
                if slide.slide_layout:
                    layout_info["layout_name"] = slide.slide_layout.name or "default"
            except Exception:
                pass

            # Detect layout type based on placeholders
            for shape in slide.shapes:
                # Check for title placeholder
                if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                    ph_type = shape.placeholder_format.type
                    if "TITLE" in str(ph_type) or "CENTER_TITLE" in str(ph_type):
                        layout_info["has_title"] = True
                        layout_info["placeholder_count"] += 1
                    elif "BODY" in str(ph_type) or "CONTENT" in str(ph_type):
                        layout_info["has_body"] = True
                        layout_info["placeholder_count"] += 1

                # Check for chart placeholder
                if shape.has_chart:
                    layout_info["has_chart"] = True
                    layout_info["placeholder_count"] += 1

                # Check for table placeholder
                if shape.has_table:
                    layout_info["has_table"] = True
                    layout_info["placeholder_count"] += 1

            # Determine layout type
            if layout_info["has_chart"]:
                layout_info["layout_type"] = "chart"
            elif layout_info["has_table"]:
                layout_info["layout_type"] = "table"
            elif layout_info["has_title"] and layout_info["has_body"]:
                if layout_info["placeholder_count"] > 2:
                    layout_info["layout_type"] = "two_column"
                else:
                    layout_info["layout_type"] = "content"
            elif layout_info["has_title"]:
                layout_info["layout_type"] = "title"
            else:
                layout_info["layout_type"] = "blank"

            layouts.append(layout_info)

        return layouts

    def _extract_placeholders(self, slide: Any) -> list[PlaceholderInfo]:
        """Extract {{variable}} placeholders from a slide with shape info.

        Args:
            slide: python-pptx slide object

        Returns:
            List of PlaceholderInfo objects with position and shape details
        """
        placeholders: list[PlaceholderInfo] = []

        for shape in slide.shapes:
            shape_id = shape.shape_id
            shape_name = shape.name

            # Determine placeholder type
            placeholder_type: str | None = None
            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                ph_type = str(shape.placeholder_format.type)
                if "TITLE" in ph_type:
                    placeholder_type = "title"
                elif "BODY" in ph_type or "CONTENT" in ph_type:
                    placeholder_type = "body"
                elif "CHART" in ph_type:
                    placeholder_type = "chart"
                elif "TABLE" in ph_type:
                    placeholder_type = "table"
                elif "PICTURE" in ph_type:
                    placeholder_type = "image"
                else:
                    placeholder_type = "other"

            # Extract text from text frames
            if hasattr(shape, "text_frame"):
                text = shape.text or ""
                if self._contains_placeholder(text):
                    # Create placeholder info for the shape
                    variables = self._extract_placeholders_from_text(text)
                    for var_name in variables:
                        placeholders.append(PlaceholderInfo(
                            variable_name=var_name,
                            shape_id=shape_id,
                            shape_name=shape_name,
                            placeholder_type=placeholder_type,
                            text=text,
                            position=(0, 0),  # Default position
                        ))

            # Also check paragraphs and runs for more granular extraction
            if hasattr(shape, "text_frame"):
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        run_text = run.text or ""
                        if self._contains_placeholder(run_text):
                            variables = self._extract_placeholders_from_text(run_text)
                            for var_name in variables:
                                placeholders.append(PlaceholderInfo(
                                    variable_name=var_name,
                                    shape_id=shape_id,
                                    shape_name=shape_name,
                                    placeholder_type=placeholder_type,
                                    text=run_text,
                                    position=(para_idx, run_idx),
                                ))

        return placeholders

    def _extract_placeholders_from_text(self, text: str) -> list[str]:
        """Extract variable names from text containing {{variable}} patterns.

        Args:
            text: Text containing {{variable_name}} patterns

        Returns:
            List of variable names found
        """
        return extract_placeholders(text)

    def _map_content_to_slides(
        self,
        content: str,
        slide_map: list[SlideStructure],
    ) -> dict[int, dict[str, Any]]:
        """Map document content to slide structure.

        Args:
            content: Document content (markdown-like text)
            slide_map: List of SlideStructure objects

        Returns:
            Dictionary mapping slide_index to content for placeholders
        """
        content_map: dict[int, dict[str, Any]] = {}

        # Parse content into sections
        sections = self._parse_content_sections(content)

        # Map sections to slides based on placeholder types
        for idx, slide_struct in enumerate(slide_map):
            slide_content: dict[str, Any] = {}

            for placeholder in slide_struct.placeholders:
                var_name = placeholder.variable_name
                ph_type = placeholder.placeholder_type

                # Find matching content for this placeholder
                mapped_content = self._find_content_for_placeholder(
                    var_name, ph_type, sections
                )

                if mapped_content:
                    slide_content[var_name] = {
                        "content": mapped_content,
                        "placeholder_type": ph_type,
                        "shape_name": placeholder.shape_name,
                    }

            if slide_content:
                content_map[idx] = slide_content

        return content_map

    def _parse_content_sections(self, content: str) -> dict[str, str]:
        """Parse document content into sections by headers.

        Args:
            content: Document content

        Returns:
            Dictionary mapping section headers/tags to content
        """
        sections: dict[str, str] = {}
        current_section = "default"
        current_content: list[str] = []

        lines = content.split("\n")
        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Check for header
            if line.startswith("# "):
                # Save previous section
                if current_content:
                    sections[current_section] = "\n".join(current_content)
                    current_content = []
                current_section = line[2:].strip()
            elif line.startswith("## "):
                # Save previous section
                if current_content:
                    sections[current_section] = "\n".join(current_content)
                    current_content = []
                current_section = line[3:].strip()
            else:
                current_content.append(line)

        # Save last section
        if current_content:
            sections[current_section] = "\n".join(current_content)

        return sections

    def _find_content_for_placeholder(
        self,
        var_name: str,
        placeholder_type: str | None,
        sections: dict[str, str],
    ) -> str | None:
        """Find content that matches a placeholder.

        Args:
            var_name: Variable name (e.g., 'title', 'executive_summary')
            placeholder_type: Type of placeholder (title, body, chart, etc.)
            sections: Parsed content sections

        Returns:
            Content to fill in, or None if not found
        """
        # Try exact match on variable name
        var_lower = var_name.lower()
        for section_key, section_content in sections.items():
            key_lower = section_key.lower()
            if var_lower in key_lower or key_lower in var_lower:
                return section_content

        # Try to match by placeholder type
        if placeholder_type == "title":
            # Return first non-default section or 'default'
            for key, value in sections.items():
                if key != "default":
                    return key
            return sections.get("default")

        if placeholder_type == "body":
            # Return longest section or 'default' section
            if "default" in sections:
                return sections["default"]
            return max(sections.values(), key=len) if sections else None

        # Return 'default' section for other types
        return sections.get("default")

    def _fill_slide_placeholders(
        self,
        slide: Any,
        content_map: dict[str, Any],
    ) -> None:
        """Replace placeholders with actual content, preserving styles.

        Args:
            slide: python-pptx slide object
            content_map: Dictionary mapping variable names to content
        """
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                self._fill_shape_text_frame(shape, content_map)

    def _fill_shape_text_frame(
        self,
        shape: Any,
        content_map: dict[str, Any],
    ) -> None:
        """Fill a shape's text frame with mapped content.

        Args:
            shape: python-pptx shape object with text_frame
            content_map: Dictionary mapping variable names to content
        """
        tf = shape.text_frame

        # First pass: check if entire text frame contains placeholders
        full_text = tf.text or ""
        if self._contains_placeholder(full_text):
            # Replace entire frame content
            new_text = self._substitute_variables(full_text, content_map)
            self._set_text_frame_preserving_style(tf, new_text)
            return

        # Second pass: process paragraph by paragraph
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run_text = run.text or ""
                if self._contains_placeholder(run_text):
                    new_text = self._substitute_variables(run_text, content_map)
                    # Preserve run properties (font, size, color)
                    run.text = new_text

    def _set_text_frame_preserving_style(
        self,
        text_frame: Any,
        new_text: str,
    ) -> None:
        """Set text frame content while preserving paragraph styles.

        Args:
            text_frame: python-pptx TextFrame object
            new_text: New text to set
        """
        # Save first paragraph's properties
        if text_frame.paragraphs:
            first_para = text_frame.paragraphs[0]
            level = first_para.level if hasattr(first_para, "level") else 0
            alignment = first_para.alignment if hasattr(first_para, "alignment") else PP_ALIGN.LEFT
        else:
            level = 0
            alignment = PP_ALIGN.LEFT

        # Clear and set new text
        text_frame.clear()

        # Add paragraphs
        lines = new_text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = line
            p.level = level
            if hasattr(p, "alignment"):
                p.alignment = alignment

    def get_presentation_structure(self, prs: Presentation) -> dict[str, Any]:
        """Get complete structure of a presentation.

        Args:
            prs: Presentation object

        Returns:
            Dictionary with complete presentation structure
        """
        slide_layouts = self._detect_slide_layouts(prs)

        slide_structures: list[SlideStructure] = []
        for idx, slide in enumerate(prs.slides):
            placeholders = self._extract_placeholders(slide)
            layout_info = slide_layouts[idx] if idx < len(slide_layouts) else {}

            # Convert placeholders to dict format
            placeholder_list = [
                {
                    "variable_name": p.variable_name,
                    "shape_id": p.shape_id,
                    "shape_name": p.shape_name,
                    "placeholder_type": p.placeholder_type,
                    "position": p.position,
                }
                for p in placeholders
            ]

            # Get shapes with placeholders
            shapes_with_placeholders = list(set([p.shape_name for p in placeholders]))

            slide_structures.append({
                "slide_index": idx,
                "layout_name": layout_info.get("layout_name", "unknown"),
                "layout_type": layout_info.get("layout_type", "unknown"),
                "placeholder_count": len(placeholders),
                "placeholders": placeholder_list,
                "shapes_with_placeholders": shapes_with_placeholders,
            })

        return {
            "slide_count": len(prs.slides),
            "slides": slide_structures,
            "layouts": slide_layouts,
        }
